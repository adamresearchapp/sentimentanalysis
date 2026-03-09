# ---------------------------------------------------------
# Media Intelligence Dashboard (dashboardV2.py)
# Hardened version aligned with databuilder.py
# ---------------------------------------------------------

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from pathlib import Path
from dataclasses import dataclass
from typing import List
import altair as alt
import pandas as pd
import streamlit as st
from wordcloud import WordCloud 
import matplotlib.pyplot as plt
import numpy as np


import os
from pathlib import Path


import altair as alt

def presschoice_dark_theme():
    return {
        "config": {
            "background": "#0E1117",

            "view": {
                "fill": "#0E1117",
                "stroke": "transparent",
                "continuousWidth": 1200,
                "continuousHeight": 650,
            },

            "axis": {
                "labelFontSize": 24,
                "titleFontSize": 24,
                "labelAngle": 0,
                "labelAlign": "center",
                "labelBaseline": "middle",
                "labelPadding": 12,
                "labelOverlap": "greedy",
                "labelLimit": 400,
                "titlePadding": 10,
                "labelColor": "#FFFFFF",
                "titleColor": "#FFFFFF",
            },

            "legend": {
                "labelFontSize": 32,
                "titleFontSize": 22,
                "labelColor": "#FFFFFF",
                "titleColor": "#FFFFFF",
            },

            "title": {
                "fontSize": 28,
                "color": "#FFFFFF",
            },
            "axisY": {
            "labelAlign": "right",
            "labelPadding": 12,
            "titlePadding": 20
        },
        }
    }


alt.themes.register("presschoice_dark", presschoice_dark_theme)
alt.themes.enable("presschoice_dark")


# Standard chart sizing constants
DEFAULT_CHART_WIDTH = 1100
DEFAULT_CHART_HEIGHT = 650
# Larger export sizes for PPTX/PNG
EXPORT_CHART_WIDTH = 1400
EXPORT_CHART_HEIGHT = 850
# Scale factor used when saving charts to disk (higher => sharper PNG)
EXPORT_SCALE_FACTOR = 4




CHART_EXPORT_DIR = Path("./powerpoint")
CHART_EXPORT_DIR.mkdir(exist_ok=True)

def save_chart(chart: alt.Chart, filename: str) -> str:
    path = CHART_EXPORT_DIR / filename
    # Save higher-resolution charts for PowerPoint export and ensure title/size preserved
    # Ensure chart has consistent dark background and readable axis/title colours for export
    try:
        chart = chart.configure_view(fill="#0E1117")
        chart = chart.configure_axis(labelColor="#FFFFFF", titleColor="#FFFFFF")
        chart = chart.configure_title(color="#FFFFFF")
        chart = chart.configure_legend(labelColor="#FFFFFF", titleColor="#FFFFFF")
        # add consistent padding so labels don't get clipped when exported
        chart = chart.properties(padding={"left": 60, "right": 20, "top": 40, "bottom": 60})
    except Exception:
        # best-effort: continue if configure not available for this chart
        pass

    chart.save(path, scale_factor=EXPORT_SCALE_FACTOR)
    return str(path)


# ---------------------------------------------------------
# PATH TO MASTER JSON
# ---------------------------------------------------------

MASTER_JSON = Path("./master.json")

# ---------------------------------------------------------
# SENTIMENT LABEL CLEANUP
# ---------------------------------------------------------

SENTIMENT_LABEL_DISPLAY = {
    "very_negative": "Very Negative",
    "negative": "Negative",
    "neutral": "Neutral",
    "positive": "Positive",
    "very_positive": "Very Positive",
}

SENTIMENT_ORDER = [
    "Very Negative",
    "Negative",
    "Neutral",
    "Positive",
    "Very Positive",
]

SENTIMENT_WEIGHTS = {
    "Very Positive": 2,
    "Positive": 1,
    "Neutral": 0,
    "Negative": -1,
    "Very Negative": -2,
}

SENTIMENT_COLORS = {
    "Very Negative": "#d73027",
    "Negative": "#fc8d59",
    "Neutral": "#808080",
    "Positive": "#66bd63",
    "Very Positive": "#1a9850",
}

# ---------------------------------------------------------
# TOPIC BUCKET MAPPING
# ---------------------------------------------------------

TOPIC_BUCKET_MAP = {
    # Performance & Strategy
    "Financial Performance & Market Position": "Performance & Strategy",
    "Strategy & Transformation": "Performance & Strategy",

    # Customer & Brand Experience
    "Customer Experience & Service Delivery": "Customer & Brand Experience",
    "Corporate Reputation & Public Perception": "Customer & Brand Experience",
    "Products & Offerings": "Customer & Brand Experience",

    # Governance, Leadership & Accountability
    "Leadership": "Governance, Leadership & Accountability",
    "Regulation & Compliance": "Governance, Leadership & Accountability",

    # Workforce, Culture & Operations
    "Workforce, Culture & Operations": "Workforce, Culture & Operations",
}

BUCKET_SHORT = {
    "Performance & Strategy": "Performance\n& Strategy",
    "Customer & Brand Experience": "Customer\n& Brand",
    "Governance, Leadership & Accountability": "Governance\n& Leadership",
    "Workforce, Culture & Operations": "Workforce\n& Operations",
}

TOPIC_THRESHOLD = 0.18
TOPIC_ALPHA_EMBED = 0.6
TOPIC_ALPHA_NLI = 0.4

# ---------------------------------------------------------
# LOAD MASTER JSON
# ---------------------------------------------------------

@st.cache_data
def load_master():
    data = json.loads(MASTER_JSON.read_text(encoding="utf-8"))

    df_sent = pd.DataFrame(data.get("sentences", []))
    df_topics = pd.DataFrame(data.get("topics", []))
    df_entities = pd.DataFrame(data.get("entities_corpus", []))
    df_linked = pd.DataFrame(data.get("entities_linked", []))
    df_ent_sent = pd.DataFrame(data.get("entity_sentiment", []))
    df_ent_time = pd.DataFrame(data.get("entity_timeline", []))

    # Fallbacks to avoid KeyError
    if "topic_name" not in df_sent.columns:
        df_sent["topic_name"] = "None"
    if "topic_score" not in df_sent.columns:
        df_sent["topic_score"] = 0.0

    if "topic_name" not in df_topics.columns and not df_topics.empty:
        df_topics["topic_name"] = df_sent["topic_name"]
    if "topic_score" not in df_topics.columns and not df_topics.empty:
        df_topics["topic_score"] = df_sent["topic_score"]

    # Always create sentiment_display
    if "label_5" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["label_5"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_label" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_label"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_class" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_class"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_category" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_category"].map(SENTIMENT_LABEL_DISPLAY)
    else:
        df_sent["sentiment_display"] = "Neutral"

    return data, df_sent, df_topics, df_entities, df_linked, df_ent_sent, df_ent_time

# ---------------------------------------------------------
# APPLY TOPIC BUCKETS
# ---------------------------------------------------------

def apply_topic_buckets(df_sent, df_topics):
    df_sent = df_sent.copy()
    df_topics = df_topics.copy()

    if "topic_name" not in df_sent.columns:
        df_sent["topic_name"] = "None"
    if "topic_name" not in df_topics.columns:
        df_topics["topic_name"] = "None"

    df_sent["topic_bucket"] = df_sent["topic_name"].map(TOPIC_BUCKET_MAP).fillna("None")
    df_topics["topic_bucket"] = df_topics["topic_name"].map(TOPIC_BUCKET_MAP).fillna("None")

    if "sentence" not in df_sent.columns:
        df_sent["sentence"] = df_sent.get("text", "")

    bucket_sizes = (
        df_sent[df_sent["topic_bucket"].ne("None")]
        .groupby("topic_bucket")["sentence"]
        .count()
        .reset_index()
        .rename(columns={"sentence": "size"})
        .sort_values("size", ascending=False)
    )

    return df_sent, df_topics, bucket_sizes

# ---------------------------------------------------------
# SECTION 2 — Bucket Analytics
# ---------------------------------------------------------

def compute_bucket_polarity(df_sent: pd.DataFrame) -> pd.DataFrame:
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return pd.DataFrame()

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    rows = []
    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]

        polarity = sum(
            SENTIMENT_WEIGHTS[row["sentiment_display"]] * row["percent"]
            for _, row in sub.iterrows()
        )

        pos = sub[sub["sentiment_display"].isin(["Positive", "Very Positive"])]["percent"].sum()
        neg = sub[sub["sentiment_display"].isin(["Negative", "Very Negative"])]["percent"].sum()
        neu = sub[sub["sentiment_display"].eq("Neutral")]["percent"].sum()

        rows.append({
            "topic_bucket": bucket,
            "bucket_short": BUCKET_SHORT.get(bucket, bucket),
            "polarity": polarity,
            "positive_percent": pos,
            "negative_percent": neg,
            "neutral_percent": neu,
        })

    return pd.DataFrame(rows)


def compute_article_bucket_polarity(df_sent: pd.DataFrame) -> pd.DataFrame:
    """Compute article-level polarity per topic bucket.

    Steps:
    - Group sentences by `article_id` and compute an article sentiment using a weighted sum of
      sentiment weights * topic_score. If the weighted sum is >0 => Positive, <0 => Negative, else Neutral.
    - Then compute percent of articles per `topic_bucket` that are Positive/Negative/Neutral.
    Returns a dataframe with article-level percentages per bucket.
    """
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty or "article_id" not in df.columns:
        return pd.DataFrame()

    # Map sentiment weights
    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS).fillna(0.0)

    # Compute article-level weighted score: sum(sentiment_weight * topic_score)
    df["weighted"] = df["sentiment_weight"] * df.get("topic_score", 1.0)

    art_scores = (
        df.groupby("article_id")["weighted"].sum().reset_index().rename(columns={"weighted": "article_score"})
    )

    # For articles with no score (shouldn't happen), fall back to majority sentiment
    # Merge back to one representative topic_bucket per article (most frequent bucket)
    rep_bucket = (
        df.groupby(["article_id", "topic_bucket"]).size().reset_index(name="count")
        .sort_values(["article_id", "count"], ascending=[True, False])
        .groupby("article_id").first().reset_index()[["article_id", "topic_bucket"]]
    )

    art = art_scores.merge(rep_bucket, on="article_id", how="left")

    # Assign article-level sentiment label
    def label_from_score(s):
        if s > 0:
            return "Positive"
        if s < 0:
            return "Negative"
        return "Neutral"

    art["article_sentiment"] = art["article_score"].apply(label_from_score)

    # Now compute percent distribution per topic_bucket
    pivot = (
        art.groupby(["topic_bucket", "article_sentiment"]).size().reset_index(name="count")
    )
    totals = (
        pivot.groupby("topic_bucket")["count"].sum().reset_index().rename(columns={"count": "total_articles"})
    )
    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total_articles"]) * 100

    rows = []
    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]
        pos = sub[sub["article_sentiment"] == "Positive"]["percent"].sum()
        neg = sub[sub["article_sentiment"] == "Negative"]["percent"].sum()
        neu = sub[sub["article_sentiment"] == "Neutral"]["percent"].sum()
        total_articles = int(sub["total_articles"].iloc[0]) if not sub.empty else 0

        rows.append({
            "topic_bucket": bucket,
            "bucket_short": BUCKET_SHORT.get(bucket, bucket),
            "article_positive_percent": pos,
            "article_negative_percent": neg,
            "article_neutral_percent": neu,
            "total_articles": total_articles,
        })

    return pd.DataFrame(rows)


def generate_bucket_summary(df_sent: pd.DataFrame) -> dict:
    summaries = {}

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return summaries

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]

        polarity = sum(
            SENTIMENT_WEIGHTS[row["sentiment_display"]] * row["percent"]
            for _, row in sub.iterrows()
        )

        pos = sub[sub["sentiment_display"].isin(["Positive", "Very Positive"])]["percent"].sum()
        neg = sub[sub["sentiment_display"].isin(["Negative", "Very Negative"])]["percent"].sum()
        neu = sub[sub["sentiment_display"].eq("Neutral")]["percent"].sum()

        if polarity > 0:
            tone = "strong positive sentiment"
            arrow = "↑"
        elif polarity < 0:
            tone = "strong negative sentiment"
            arrow = "↓"
        else:
            tone = "mixed or neutral sentiment"
            arrow = "→"

        summaries[bucket] = (
            f"{arrow} {bucket} shows {tone}. "
            f"Positive coverage: {pos:.1f}%. "
            f"Negative coverage: {neg:.1f}%. "
            f"Neutral coverage: {neu:.1f}%. "
            f"Polarity score: {polarity:.1f}."
        )

    return summaries


def get_sentiment_drivers(df_sent: pd.DataFrame, top_n: int = 5) -> dict:
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return {}

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
    df["driver_score"] = df["sentiment_weight"] * df["topic_score"]

    drivers = {}

    for bucket in df["topic_bucket"].unique():
        sub = df[df["topic_bucket"] == bucket]

        pos = (
            sub[sub["sentiment_weight"] > 0]
            .sort_values("driver_score", ascending=False)
            .head(top_n)
        )
        neg = (
            sub[sub["sentiment_weight"] < 0]
            .sort_values("driver_score")
            .head(top_n)
        )

        drivers[bucket] = {
            "positive": pos[["sentence", "sentiment_display", "topic_score", "driver_score"]],
            "negative": neg[["sentence", "sentiment_display", "topic_score", "driver_score"]],
        }

    return drivers


def get_global_sentiment_drivers(df_sent: pd.DataFrame, top_n: int = 3):
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()

    if df.empty:
        return pd.DataFrame(), pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
    df["driver_score"] = df["sentiment_weight"] * df["topic_score"]

    pos = (
        df[df["sentiment_weight"] > 0]
        .sort_values("driver_score", ascending=False)
        .head(top_n)
    )
    neg = (
        df[df["sentiment_weight"] < 0]
        .sort_values("driver_score")
        .head(top_n)
    )

    return pos, neg

# ---------------------------------------------------------
# SECTION 3 — Visuals (Heatmaps, Drift, Purity, Balance Bubbles)
# ---------------------------------------------------------

def compute_article_sentiment(df_sent: pd.DataFrame) -> pd.DataFrame:
    """
    Aggregate sentence-level sentiment into a per-article tone.
    Returns one row per article with an overall_tone label.
    """
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "article_id" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent.copy()
    df = df[df["article_id"].notna()]
    if df.empty:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)

    agg = (
        df.groupby("article_id")
        .agg(
            avg_weight=("sentiment_weight", "mean"),
            sentence_count=("sentence", "size"),
        )
        .reset_index()
    )

    def classify_tone(weight: float) -> str:
        if weight > 0.1:
            return "Positive"
        if weight < -0.1:
            return "Negative"
        return "Neutral"

    agg["overall_tone"] = agg["avg_weight"].apply(classify_tone)
    return agg


def build_sentence_distribution_chart(df_sent: pd.DataFrame):
    if df_sent is None or df_sent.empty:
        return None

    df = (
        df_sent.groupby("sentiment_display")
        .size()
        .reset_index(name="count")
    )

    df["sentiment_display"] = pd.Categorical(
        df["sentiment_display"],
        categories=SENTIMENT_ORDER,
        ordered=True,
    )
    df = df.sort_values("sentiment_display")

    color_scale = alt.Scale(
        domain=SENTIMENT_ORDER,
        range=[SENTIMENT_COLORS[s] for s in SENTIMENT_ORDER],
    )

    chart = (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(
                    title="Sentiment",
                ),
            ),
            y=alt.Y("count:Q", title="Number of sentences"),
            color=alt.Color(
                "sentiment_display:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=["sentiment_display", "count"],
        )
    ).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)

    return chart


def build_article_tone_chart(df_article_sent: pd.DataFrame):
    if df_article_sent is None or df_article_sent.empty:
        return None

    df = (
        df_article_sent.groupby("overall_tone")
        .size()
        .reset_index(name="count")
    )

    tone_order = ["Negative", "Neutral", "Positive"]
    df["overall_tone"] = pd.Categorical(
        df["overall_tone"],
        categories=tone_order,
        ordered=True,
    )
    df = df.sort_values("overall_tone")

    tone_colors = {
        "Negative": SENTIMENT_COLORS["Negative"],
        "Neutral": SENTIMENT_COLORS["Neutral"],
        "Positive": SENTIMENT_COLORS["Positive"],
    }

    color_scale = alt.Scale(
        domain=tone_order,
        range=[tone_colors[t] for t in tone_order],
    )

    chart = (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X(
                "overall_tone:N",
                sort=tone_order,
                axis=alt.Axis(
                    title="Article tone",
                ),
            ),
            y=alt.Y("count:Q", title="Number of articles"),
            color=alt.Color(
                "overall_tone:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=["overall_tone", "count"],
        )
    ).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)

    return chart


def build_bucket_balance_table(df_sent: pd.DataFrame) -> pd.DataFrame:
    """
    Per topic bucket:
    - net_balance = #positive sentences - #negative sentences  (x-axis)
    - avg_intensity = average sentiment_weight (y-axis)
    - total_count = total sentences in bucket (bubble size)
    """
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "topic_bucket" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)

    # Positive/negative masks
    is_pos = df["sentiment_display"].isin(["Positive", "Very Positive"])
    is_neg = df["sentiment_display"].isin(["Negative", "Very Negative"])

    agg = (
        df.groupby("topic_bucket")
        .agg(
            positive_count=("sentiment_display", lambda s: int(is_pos.loc[s.index].sum())),
            negative_count=("sentiment_display", lambda s: int(is_neg.loc[s.index].sum())),
            total_count=("sentiment_display", "size"),
            avg_intensity=("sentiment_weight", "mean"),
        )
        .reset_index()
    )

    agg["net_balance"] = agg["positive_count"] - agg["negative_count"]
    agg["bucket_short"] = agg["topic_bucket"].map(BUCKET_SHORT)

    return agg


def bucket_balance_bubble(df_sent: pd.DataFrame):
    """
    Bubble chart where:
    - x-axis: net positive minus negative sentences (cumulative balance)
    - y-axis: average sentiment intensity (-2 to +2)
    - bubble size: total number of sentences in the bucket
    - label: topic bucket name inside the bubble
    """
    table = build_bucket_balance_table(df_sent)
    if table is None or table.empty:
        return None

    # Make sure we have a stable order by net balance
    table = table.sort_values("net_balance")

    base = alt.Chart(table).encode(
        x=alt.X(
            "net_balance:Q",
            axis=alt.Axis(
                title="Net positive minus negative sentences",
            ),
            scale=alt.Scale(nice=True, padding=20),
        ),
        y=alt.Y(
            "avg_intensity:Q",
            axis=alt.Axis(
                title="Average sentiment intensity",
            ),
            scale=alt.Scale(nice=True, padding=10),
        ),
    )

    bubbles = base.mark_circle(opacity=0.7, stroke="black", strokeWidth=0.5).encode(
        size=alt.Size(
            "total_count:Q",
            title="Total sentences in bucket",
            scale=alt.Scale(range=[200, 2000]),
        ),
        color=alt.Color(
            "avg_intensity:Q",
            title="Average intensity",
            scale=alt.Scale(scheme="redyellowgreen"),
            legend=None,
        ),
        tooltip=[
            "topic_bucket",
            "net_balance",
            "avg_intensity",
            "positive_count",
            "negative_count",
            "total_count",
        ],
    )

    labels = (
        base.mark_text(
            baseline="middle",
            align="center",
            fontSize=11,
            color="white",        # white fill
            # stroke="black",       # thin black outline
            strokeWidth=0.8,
        ).encode(
            text=alt.Text("bucket_short:N", title=None),
        )
    )

    # Reference lines at x=0 and y=0
    vline = alt.Chart(table).mark_rule(color="#555555", strokeWidth=2).encode(x=alt.datum(0))
    hline = alt.Chart(table).mark_rule(color="#555555", strokeWidth=2).encode(y=alt.datum(0))

    return (vline + hline + bubbles + labels).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)
def bucket_sentiment_heatmap(df_sent: pd.DataFrame):
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return None

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100
    pivot["bucket_short"] = pivot["topic_bucket"].map(BUCKET_SHORT)

    chart = (
        alt.Chart(pivot)
        .mark_rect()
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(
                    title="Sentiment",
                ),
            ),
            y=alt.Y(
                "bucket_short:N",
                axis=alt.Axis(
                    title="Topic Bucket",
                ),
            ),
            color=alt.Color(
                "percent:Q",
                scale=alt.Scale(scheme="blues"),
                title="Percent of Bucket",
            ),
            tooltip=[
                "topic_bucket",
                "sentiment_display",
                "percent",
                "count",
                "total",
            ],
        ).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)
    )

    return chart


def bucket_sentiment_bubble(df_sent: pd.DataFrame):
    """
    Bubble view of sentiment by topic bucket:
    x = sentiment, y = bucket, size/color = share of bucket.
    """
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return None

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100
    pivot["bucket_short"] = pivot["topic_bucket"].map(BUCKET_SHORT)

    color_scale = alt.Scale(
        domain=SENTIMENT_ORDER,
        range=[SENTIMENT_COLORS[s] for s in SENTIMENT_ORDER],
    )

    chart = (
        alt.Chart(pivot)
        .mark_circle(opacity=0.85, stroke="black", strokeWidth=0.2)
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(
                    title="Sentiment",
                ),
            ),
            y=alt.Y(
                "bucket_short:N",
                axis=alt.Axis(
                    title="Topic Bucket",
                ),
            ),
            size=alt.Size(
                "percent:Q",
                title="Percent of Bucket",
                scale=alt.Scale(range=[50, 1000]),
            ),
            color=alt.Color(
                "sentiment_display:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=[
                "topic_bucket",
                "sentiment_display",
                "percent",
                "count",
                "total",
            ],
        ).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)
    )

    return chart


def topic_drift_heatmap(df_sent: pd.DataFrame):
    if "topic_bucket" not in df_sent.columns:
        return None

    df = df_sent.copy()
    df = df[df["topic_bucket"].ne("None")]
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return None

    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    df["count"] = 1

    pivot = (
        df.groupby(["topic_name", "bucket_short"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_name")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_name")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    chart = (
        alt.Chart(pivot)
        .mark_rect()
        .encode(
            x=alt.X(
                "bucket_short:N",
                axis=alt.Axis(
                    title="Bucket",
                ),
            ),
            y=alt.Y(
                "topic_name:N",
                axis=alt.Axis(
                    title="Fine Topic",
                ),
            ),
            color=alt.Color(
                "percent:Q",
                scale=alt.Scale(scheme="greens"),
                title="Percent of Topic",
            ),
            tooltip=[
                "topic_name",
                "bucket_short",
                "percent",
                "count",
                "total",
            ],
        ).properties(width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT)
    )

    return chart


def compute_topic_purity(df_sent: pd.DataFrame) -> pd.DataFrame:
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "topic_bucket" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return pd.DataFrame()

    df["count"] = 1

    pivot = (
        df.groupby(["topic_name", "topic_bucket"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_name")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_name")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    purity = (
        pivot.groupby("topic_name")["percent"]
        .max()
        .reset_index()
        .rename(columns={"percent": "purity"})
    )

    return purity.sort_values("purity", ascending=False)

# ---------------------------------------------------------
# SECTION 4 — Executive Summary Page
# ---------------------------------------------------------

def build_executive_summary(df_polarity: pd.DataFrame) -> str:
    if df_polarity.empty:
        return "No sentiment data available for an executive summary."

    lines = []

    for _, row in df_polarity.sort_values("polarity", ascending=False).iterrows():
        bucket = row["topic_bucket"]
        pol = row["polarity"]
        pos = row["positive_percent"]
        neg = row["negative_percent"]
        neu = row["neutral_percent"]

        if pol > 0:
            direction = "is a positive area with favourable coverage."
        elif pol < 0:
            direction = "is a negative area with critical coverage."
        else:
            direction = "shows a mixed or neutral sentiment profile."

        lines.append(
            f"{bucket} {direction} "
            f"Approx. {pos:.0f}% positive, {neg:.0f}% negative, and {neu:.0f}% neutral."
        )

    return " ".join(lines)


def render_executive_summary_page(df_sent: pd.DataFrame):
    st.header("Executive Summary")

    # Overall tone: sentences vs articles
    df_article_sent = compute_article_sentiment(df_sent)

    st.markdown("### Overall Sentiment — Sentences vs Articles")
    col_a, col_b = st.columns(2)

    with col_a:
        st.markdown("Sentence Sentiment Distribution")
        chart_sent = build_sentence_distribution_chart(df_sent)
        if chart_sent is not None:
            st.altair_chart(chart_sent, width="stretch")
        else:
            st.write("No sentence-level sentiment data available.")

    with col_b:
        st.markdown("Article Tone Distribution")
        chart_articles = build_article_tone_chart(df_article_sent)
        if chart_articles is not None:
            st.altair_chart(chart_articles, width="stretch")
        else:
            st.write("No article-level sentiment data available.")

    df_polarity = compute_bucket_polarity(df_sent)
    summary_text = build_executive_summary(df_polarity)

    st.markdown("### Narrative Overview")
    st.write(summary_text)

    overall = compute_overall_score(df_sent)
    st.metric("Overall Sentiment Score", f"{overall:.1f} / 100")

    st.markdown("### Polarity by Bucket")

    if not df_polarity.empty:
        chart_pol = (
            alt.Chart(df_polarity)
            .mark_bar()
            .encode(
                x=alt.X(
                    "bucket_short:N",
                    sort="-y",
                    axis=alt.Axis(
                        title="Bucket",
                    ),
                ),
                y=alt.Y("polarity:Q", title="Polarity Score"),
                color=alt.condition(
                    alt.datum.polarity > 0,
                    alt.value("#2ca02c"),
                    alt.value("#d62728"),
                ),
                tooltip=[
                    "topic_bucket",
                    "polarity",
                    "positive_percent",
                    "negative_percent",
                    "neutral_percent",
                ],
            )
        )
        st.altair_chart(chart_pol, width="stretch")
    else:
        st.write("No polarity data available.")

    st.markdown("### Top Sentiment Drivers (Global)")

    pos, neg = get_global_sentiment_drivers(df_sent, top_n=3)
    col1, col2 = st.columns(2)

    with col1:
        st.markdown("Top Positive Drivers")
        if not pos.empty:
            st.dataframe(
                pos[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "driver_score",
                    ]
                ]
            )
        else:
            st.write("No positive drivers available.")

    with col2:
        st.markdown("Top Negative Drivers")
        if not neg.empty:
            st.dataframe(
                neg[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "driver_score",
                    ]
                ]
            )
        else:
            st.write("No negative drivers available.")

# ---------------------------------------------------------
# SECTION 5 — Topic Buckets Page
# ---------------------------------------------------------

def render_topic_buckets_page(df_sent: pd.DataFrame, df_topics: pd.DataFrame, bucket_sizes: pd.DataFrame):
    st.header("High‑Level Topic Buckets")

    if "topic_bucket" not in df_sent.columns:
        df_sent, df_topics, bucket_sizes = apply_topic_buckets(df_sent, df_topics)

    df_sent_b = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df_sent_b = df_sent_b[df_sent_b["topic_bucket"].ne("Other")]

    st.markdown("### Bucket Sizes")
    if bucket_sizes is not None and not bucket_sizes.empty:
        st.dataframe(bucket_sizes[bucket_sizes["topic_bucket"].ne("None")])
    else:
        st.write("No bucket size data available.")

    df_polarity = compute_bucket_polarity(df_sent_b)

    st.markdown("### Bucket Polarity Ranking")
    if not df_polarity.empty:
        st.dataframe(df_polarity.sort_values("polarity", ascending=False))
    else:
        st.write("No polarity data available.")

    st.markdown("### Polarity by Bucket")

    if not df_polarity.empty:
        chart_pol = (
            alt.Chart(df_polarity)
            .mark_bar()
            .encode(
                x=alt.X(
                    "bucket_short:N",
                    sort="-y",
                    axis=alt.Axis(
                        title="Bucket",
                    ),
                ),
                y=alt.Y("polarity:Q", title="Polarity Score"),
                color=alt.condition(
                    alt.datum.polarity > 0,
                    alt.value("#2ca02c"),
                    alt.value("#d62728"),
                ),
                tooltip=[
                    "topic_bucket",
                    "polarity",
                    "positive_percent",
                    "negative_percent",
                    "neutral_percent",
                ],
            )
        )
        st.altair_chart(chart_pol, width="stretch")
    else:
        st.write("No polarity scores to display.")

    st.markdown("### Bucket × Sentiment View")

    view_type = st.radio(
        "Choose sentiment view",
        ["Heatmap", "Bubble chart"],
        horizontal=True,
    )

    if view_type == "Heatmap":
        heatmap = bucket_sentiment_heatmap(df_sent_b)
        if heatmap is not None:
            st.altair_chart(heatmap, width="stretch")
        else:
            st.write("No data available for bucket heatmap.")
    else:
        bubble = bucket_sentiment_bubble(df_sent_b)
        if bubble is not None:
            st.altair_chart(bubble, width="stretch")
        else:
            st.write("No data available for bucket bubble view.")

    st.markdown("### Net Balance vs Intensity (Bubble Map)")

    balance_chart = bucket_balance_bubble(df_sent_b)
    if balance_chart is not None:
        st.altair_chart(balance_chart, width="stretch")
    else:
        st.write("No data available for balance bubble map.")

    st.markdown("### Topic Drift (Fine Topics → Buckets)")

    drift_chart = topic_drift_heatmap(df_sent_b)
    if drift_chart is not None:
        st.altair_chart(drift_chart, use_container_width=True)
    else:
        st.write("No data available for topic drift.")

    st.markdown("### Topic Purity Scores")

    purity = compute_topic_purity(df_sent_b)
    if not purity.empty:
        st.dataframe(purity)
    else:
        st.write("No purity data available.")

    st.markdown("### Top Sentiment Drivers (by Bucket)")

    drivers = get_sentiment_drivers(df_sent_b)
    for bucket, d in drivers.items():
        st.markdown(f"#### {bucket}")
        col1, col2 = st.columns(2)

        with col1:
            st.markdown("Top Positive Drivers")
            if not d["positive"].empty:
                st.dataframe(d["positive"])
            else:
                st.write("No positive drivers.")

        with col2:
            st.markdown("Top Negative Drivers")
            if not d["negative"].empty:
                st.dataframe(d["negative"])
            else:
                st.write("No negative drivers.")

    st.markdown("### Narrative Summaries")

    summaries = generate_bucket_summary(df_sent_b)
    if summaries:
        for bucket, text in summaries.items():
            st.markdown(f"{bucket}")
            st.write(text)
    else:
        st.write("No summaries available.")

# ---------------------------------------------------------
# SECTION 6 — Overview, Topic Explorer, Sentence Inspector
# ---------------------------------------------------------

# --- WORD CLOUD UTILITIES ---


WORDCLOUD_EXPORT_PATH = CHART_EXPORT_DIR / "wordcloud.png"

def sentiment_to_rgb(sentiment, topic_score):
    """Return a discrete RGB colour: green for positive, red for negative, grey for neutral.
    Neutral is defined as topic_score between -5 and +5 (inclusive). If topic_score is missing,
    fall back to sentiment label (Positive/Negative) to choose colour.
    """
    # Discrete colour choices (no gradient)
    GREEN = (0, 153, 76)
    RED = (204, 0, 0)
    NEUTRAL = (150, 150, 150)

    try:
        val = float(topic_score)
    except Exception:
        val = None

    if val is not None:
        if val > 5:
            return GREEN
        if val < -5:
            return RED
        return NEUTRAL

    # Fallback to sentiment label when numeric score unavailable
    s = (sentiment or "").lower()
    if "positive" in s:
        return GREEN
    if "negative" in s:
        return RED
    return NEUTRAL



def build_sentiment_wordcloud_data(df_sent, search_terms):
    # Parse comma-separated terms
    terms = [t.strip().lower() for t in search_terms.split(",") if t.strip()]
    if not terms:
        return {}, {}

    # Filter sentences containing any of the terms
    mask = df_sent["sentence"].str.lower().apply(
        lambda s: any(t in s for t in terms)
    )
    df = df_sent[mask].copy()

    if df.empty:
        return {}, {}

    # Frequency per term
    freq = {t: df["sentence"].str.lower().str.count(t).sum() for t in terms}

    # Compute min/max topic scores for normalisation
    all_scores = df["topic_score"].tolist()
    min_score = min(all_scores)
    max_score = max(all_scores)

    # Build colour map
    colors = {}
    for t in terms:
        subset = df[df["sentence"].str.lower().str.contains(t)]
        if subset.empty:
            colors[t] = (180, 180, 180)
            continue
        # Frequency-driven size remains the same (freq)
        avg_topic = subset["topic_score"].mean()

        # Compute simple mention balance: count positive vs negative sentences for the term
        pos_count = subset[subset["sentiment_display"].str.contains("Positive", na=False)]["sentiment_display"].count()
        neg_count = subset[subset["sentiment_display"].str.contains("Negative", na=False)]["sentiment_display"].count()

        # Asymmetric thresholds: positive needs +1, negative needs -5
        POS_THRESHOLD = 3
        NEG_THRESHOLD = 3
        if pos_count - neg_count >= POS_THRESHOLD:
            # positive-dominant
            colors[t] = (0, 153, 76)
        elif neg_count - pos_count >= NEG_THRESHOLD:
            # negative-dominant
            colors[t] = (204, 0, 0)
        else:
            # neutral (within thresholds)
            colors[t] = (150, 150, 150)

    return freq, colors




def make_color_func(color_map):
    def color_func(word, *args, **kwargs):
        rgb = color_map.get(word, (180, 180, 180))
        return f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
    return color_func




def render_overview_page(df_sent, df_topics, df_articles):
    st.header("Overview")

    col1, col2, col3 = st.columns(3)

    with col1:
        st.metric("Total Sentences", len(df_sent))

    with col2:
        st.metric("Unique Topics", df_topics["topic_name"].nunique() if not df_topics.empty else 0)

    with col3:
        st.metric("Total Articles", len(df_articles) if df_articles is not None else 0)

    st.markdown("### Sample of Sentences")
    st.dataframe(df_sent.head(20))


def render_topic_explorer_page(df_sent):
    st.header("Topic Explorer")

    if "topic_name" not in df_sent.columns or df_sent["topic_name"].nunique() == 0:
        st.write("No topic data available.")
        return

    topic_list = sorted(df_sent["topic_name"].dropna().unique())
    topic = st.selectbox("Select a topic", topic_list)

    df_t = df_sent[df_sent["topic_name"] == topic].copy()

    st.markdown(f"### Sentences for: {topic}")
    st.dataframe(
        df_t[
            [
                "sentence",
                "topic_name",
                "topic_bucket",
                "sentiment_display",
                "topic_score",
            ]
        ]
    )

def render_sentence_inspector_page(df_sent):
    st.header("Sentence Inspector")

    query = st.text_input("Search text in sentences (single-term)")
    multi_terms = st.text_input("Word Cloud Terms (comma-separated)")

    # --- Table Search ---
    if query:
        df_q = df_sent[
            df_sent["sentence"].str.contains(query, case=False, na=False)
        ].copy()

        st.write(f"Found {len(df_q)} matching sentences.")
        st.dataframe(
            df_q[
                [
                    "sentence",
                    "topic_name",
                    "topic_bucket",
                    "sentiment_display",
                    "topic_score",
                ]
            ]
        )
    else:
        st.write("Enter text above to search through all sentences.")

    # --- WORD CLOUD ---
    if multi_terms:
        freq, colors = build_sentiment_wordcloud_data(df_sent, multi_terms)
        if freq:
            wc = WordCloud(
                width=1200,
                height=600,
                background_color="#0E1117",
                prefer_horizontal=1.0
            ).generate_from_frequencies(freq)

            wc.recolor(color_func=make_color_func(colors))

            fig, ax = plt.subplots(figsize=(12, 6))
            ax.imshow(wc, interpolation="bilinear")
            ax.axis("off")

            # Save for PPTX export
            wc.to_file(str(WORDCLOUD_EXPORT_PATH))

            st.pyplot(fig)
        else:
            st.info("No matching sentences for word cloud.")


# ---------------------------------------------------------
# SECTION 7 — Entity Explorer, Article Browser, Search
# ---------------------------------------------------------

def render_entity_explorer_page(df_entities, df_sent):
    st.header("Entity Explorer")

    if df_entities.empty:
        st.write("No entity data available.")
        return

    entity_list = sorted(df_entities["entity_canonical"].dropna().unique())
    entity = st.selectbox("Select an entity", entity_list)

    df_e = df_entities[df_entities["entity_canonical"] == entity].copy()
    st.markdown(f"### Mentions of: {entity}")
    st.dataframe(df_e)

    if "sentence_id" in df_e.columns and "global_index" in df_sent.columns:
        linked = df_sent[df_sent["global_index"].isin(df_e["sentence_id"])]
        if not linked.empty:
            st.markdown("### Sentences mentioning this entity")
            st.dataframe(
                linked[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "topic_score",
                    ]
                ]
            )


def render_article_browser_page(df_articles, df_sent):
    st.header("Article Browser")

    if df_articles.empty:
        st.write("No article data available.")
        return

    article_ids = df_articles["article_id"].tolist()
    article_id = st.selectbox("Select an article ID", article_ids)

    df_a = df_articles[df_articles["article_id"] == article_id].copy()
    st.markdown("### Article Metadata")
    st.dataframe(df_a)

    df_s = df_sent[df_sent["article_id"] == article_id].copy()
    st.markdown("### Sentences in this Article")
    st.dataframe(
        df_s[
            [
                "sentence",
                "topic_name",
                "topic_bucket",
                "sentiment_display",
                "topic_score",
            ]
        ]
    )


def render_search_page(df_sent, df_articles):
    st.header("Search")

    query = st.text_input("Search across sentences and articles")

    if not query:
        st.write("Enter a search term to begin.")
        return

    df_s = df_sent[df_sent["sentence"].str.contains(query, case=False, na=False)]
    st.markdown("### Sentence Matches")
    st.write(f"{len(df_s)} matches")
    st.dataframe(
        df_s[
            [
                "sentence",
                "topic_name",
                "topic_bucket",
                "sentiment_display",
                "topic_score",
            ]
        ]
    )

    if not df_articles.empty and "title" in df_articles.columns:
        df_a = df_articles[df_articles["title"].str.contains(query, case=False, na=False)]
        st.markdown("### Article Title Matches")
        st.write(f"{len(df_a)} matches")
        st.dataframe(df_a)

# ---------------------------------------------------------
# SECTION 8 — powerpoint export and narrative export
# ---------------------------------------------------------
def build_narrative(df_sent):
    df_polarity = compute_bucket_polarity(df_sent)
    pos, neg = get_global_sentiment_drivers(df_sent, top_n=5)

    def fmt_driver(df):
        if df.empty:
            return "No drivers."
        lines = []
        for _, r in df.iterrows():
            lines.append(f"“{r['sentence']}” → {r['topic_bucket']} / {r['topic_name']} (score {r['driver_score']:.2f}).")
        return "\n".join(lines)

    pos_text = fmt_driver(pos)
    neg_text = fmt_driver(neg)

    bucket_lines = []
    for _, row in df_polarity.iterrows():
        pol = row["polarity"]
        if pol > 5:
            arrow = "↑"; tone = "strong positive sentiment"
        elif pol < -5:
            arrow = "↓"; tone = "strong negative sentiment"
        else:
            arrow = "→"; tone = "mixed or neutral sentiment"

        bucket_lines.append(
            f"{arrow} {row['topic_bucket']} shows {tone}. "
            f"Positive: {row['positive_percent']:.1f}%. Negative: {row['negative_percent']:.1f}%. Neutral: {row['neutral_percent']:.1f}%. Polarity score: {pol:.1f}."
        )

    global_lines = []
    for _, row in df_polarity.iterrows():
        pol = row["polarity"]
        if pol > 5:
            tone = "a positive area with favourable coverage"
        elif pol < -5:
            tone = "a negative area with critical coverage"
        else:
            tone = "a mixed or neutral sentiment profile"
        global_lines.append(
            f"{row['topic_bucket']} is {tone}. Approx. "
            f"{row['positive_percent']:.0f}% positive, "
            f"{row['negative_percent']:.0f}% negative, "
            f"{row['neutral_percent']:.0f}% neutral."
        )

    chart_text = """
Sentence-Level Sentiment Distribution
Shows a predominantly neutral media environment with positive sentiment concentrated in performance-related topics.

Bucket Polarity
Highlights which thematic areas carry favourable or unfavourable sentiment pressure.

Sentiment Composition Bubble Chart
Bubble size reflects coverage volume; colour and position show tonal skew.

Net Balance vs Intensity
Identifies reputational strengths (upper-right quadrant) and risk zones (lower-left quadrant).
"""

    return f"""
# Top Positive Drivers
{pos_text}

# Top Negative Drivers
{neg_text}

# Narrative Summaries by Bucket
{'\n'.join(bucket_lines)}

# Global Narrative Overview
{' '.join(global_lines)}

# Chart Interpretation Summaries
{chart_text}
"""



@dataclass
class SlideSpec:
    title: str
    body_text: str
    image_path: str | None = None  # PNG path for chart (optional)
    caption: str | None = None

def compute_global_drivers(df_sent: pd.DataFrame, top_n: int = 5) -> pd.DataFrame:
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
    df["overall_driver_score"] = df["sentiment_weight"] * df["topic_score"]

    df = df.sort_values("overall_driver_score", ascending=False)
    return df.head(top_n)

def build_driver_narrative(df_sent: pd.DataFrame, top_n: int = 3) -> str:
    """Builds Narrative Summaries and top drivers per bucket.

    - Produces the bucket-level narrative lines (polarity, % pos/neg/neu) consistent with
      `compute_bucket_polarity` output.
    - Appends the top positive and negative driver (one each) per bucket where available.
    """
    if df_sent is None or df_sent.empty:
        return "No clear global sentiment drivers were identified in this period."

    # Sentence-level polarity summary
    df_polarity = compute_bucket_polarity(df_sent)

    # Top drivers by bucket (one positive, one negative)
    drivers = get_sentiment_drivers(df_sent, top_n=top_n)

    lines = []
    # Keep original bucket ordering if available, else use df_polarity order
    buckets = df_polarity['topic_bucket'].tolist() if not df_polarity.empty else []

    for bucket in buckets:
        row = df_polarity[df_polarity['topic_bucket'] == bucket].iloc[0]
        pol = row['polarity']
        pos = row['positive_percent']
        neg = row['negative_percent']
        neu = row['neutral_percent']

        if pol > 0:
            arrow = '↑'
            tone = 'shows strong positive sentiment'
        elif pol < 0:
            arrow = '↓'
            tone = 'shows strong negative sentiment'
        else:
            arrow = '→'
            tone = 'shows mixed or neutral sentiment'

        lines.append(f"{arrow} {bucket}\n{bucket} {tone}. Positive coverage: {pos:.1f}%. Negative coverage: {neg:.1f}%. Neutral coverage: {neu:.1f}%. Polarity score: {pol:.1f}.")

        # Add top positive and negative driver quotes where present
        bdrivers = drivers.get(bucket, {})
        pos_df = bdrivers.get('positive') if bdrivers else None
        neg_df = bdrivers.get('negative') if bdrivers else None

        if pos_df is not None and not pos_df.empty:
            s = pos_df.iloc[0]['sentence']
            s = (s[:197] + '…') if len(s) > 200 else s
            score = pos_df.iloc[0]['driver_score']
            lines.append(f"Top positive driver: “{s}” (score {score:.2f})")

        if neg_df is not None and not neg_df.empty:
            s = neg_df.iloc[0]['sentence']
            s = (s[:197] + '…') if len(s) > 200 else s
            score = neg_df.iloc[0]['driver_score']
            lines.append(f"Top negative driver: “{s}” (score {score:.2f})")

        # separator line
        lines.append("")

    return "\n".join(lines)


def build_media_reputation_action_plan(df_polarity=None, overall_score=None, df_sent=None) -> str:
    """
    Executive-grade Media & Reputation Action Plan.
    Advisory tone suitable for media teams.
    One bullet per bucket. No repetition. No magic numbers.
    """

    # === 1. Initialise placeholders ===
    perf_net = None
    gov_net = None
    overall = f"{overall_score:.1f}" if overall_score is not None else None

    # === 2. Extract headline metrics from df_polarity ===
    try:
        if df_polarity is not None and not df_polarity.empty:

            # Performance & Strategy
            perf_row = df_polarity[
                df_polarity["topic_bucket"]
                .str.contains("Performance|Strategy", case=False, na=False)
            ]
            if not perf_row.empty and "polarity" in perf_row.columns:
                perf_net = f"{perf_row['polarity'].iloc[0]:+.1f}%"

            # Governance
            gov_row = df_polarity[
                df_polarity["topic_bucket"]
                .str.contains("Governance|Leadership", case=False, na=False)
            ]
            if not gov_row.empty and "polarity" in gov_row.columns:
                gov_net = f"{gov_row['polarity'].iloc[0]:+.1f}%"

    except Exception:
        pass

    # === 3. Build headline line (hide missing metrics) ===
    headline_parts = []
    if perf_net:
        headline_parts.append(f"Performance & Strategy: net {perf_net}")
    if gov_net:
        headline_parts.append(f"Governance/Leadership: net {gov_net}")
    if overall:
        headline_parts.append(f"Overall: {overall}/100")

    parts = [" — ".join(headline_parts), ""]

    # === 4. Extract top global drivers (max 3) ===
    notable = []
    try:
        if df_sent is not None and not df_sent.empty:
            global_drivers = compute_global_drivers(df_sent, top_n=3)
            for _, g in global_drivers.iterrows():
                s = (g.get("sentence") or "").strip()
                s = (s[:197] + "…") if len(s) > 200 else s
                notable.append(f"“{s}” — {g.get('topic_bucket','')}/{g.get('topic_name','')}.")
    except Exception:
        pass

    if notable:
        parts.append("**Notable Global Drivers**")
        parts.extend([f"• {n}" for n in notable])
        parts.append("")

    # === 5. Build bucket-level Amplify/Mitigate sections (advisory tone) ===
    amplify = {}
    mitigate = {}

    pos_templates = {
        "Performance": "There may be an opportunity to reinforce financial and execution proof points (e.g., CFO briefings, one‑pagers, case studies).",
        "Customer": "Positive customer‑experience signals could be highlighted through testimonials, adoption stats or targeted trade outreach.",
        "Governance": "Leadership and governance narratives may be supported through roadmap updates, milestone communication or Q&A formats.",
        "Workforce": "Culture and talent themes could be reinforced through leader profiles, retention metrics or employee stories.",
    }

    neg_templates = {
        "Performance": "Short‑term market or performance movements may benefit from contextualisation using KPIs and forward‑looking commentary.",
        "Customer": "Customer‑related issues could be eased through rapid‑response comms, case‑level remediation and success stories.",
        "Governance": "Governance concerns may be mitigated through clearer timelines, independent reviews or increased leadership visibility.",
        "Workforce": "Workforce narratives may be supported by highlighting cultural programmes, retention initiatives or training investments.",
    }

    try:
        if df_sent is not None and not df_sent.empty:

            df = df_sent[df_sent["topic_bucket"].notna()].copy()
            df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
            df["driver_score"] = df["sentiment_weight"] * df.get("topic_score", 1.0)

            # Group by bucket and pick top positive + top negative driver
            for bucket, group in df.groupby("topic_bucket"):

                # Top positive driver → Amplify Opportunities
                pos = group[group["sentiment_weight"] > 0]
                if not pos.empty:
                    r = pos.sort_values("driver_score", ascending=False).iloc[0]
                    sent = (r.get("sentence") or "").strip()
                    short = (sent[:197] + "…") if len(sent) > 200 else sent
                    template = next(
                        (t for k, t in pos_templates.items() if k.lower() in bucket.lower()),
                        "There may be an opportunity to reinforce positive narratives through spokesperson commentary or case studies."
                    )
                    amplify[bucket] = f"• {template} Example signal → “{short}”."

                # Top negative driver → Mitigation Considerations
                neg = group[group["sentiment_weight"] < 0]
                if not neg.empty:
                    r = neg.sort_values("driver_score").iloc[0]
                    sent = (r.get("sentence") or "").strip()
                    short = (sent[:197] + "…") if len(sent) > 200 else sent
                    template = next(
                        (t for k, t in neg_templates.items() if k.lower() in bucket.lower()),
                        "Narratives may benefit from proactive clarification, transparent facts or targeted owned/earned content."
                    )
                    mitigate[bucket] = f"• {template} Example signal → “{short}”."

    except Exception:
        pass

    # === 6. Sort buckets by importance ===
    bucket_order = ["Performance", "Governance", "Customer", "Workforce"]
    def sort_key(b):
        for i, k in enumerate(bucket_order):
            if k.lower() in b.lower():
                return i
        return 999

    # === 7. Output Amplify section ===
    parts.append("**Potential Opportunities to Reinforce**")
    if amplify:
        for bucket in sorted(amplify.keys(), key=sort_key):
            parts.append(amplify[bucket])
    else:
        parts.append("• There may be opportunities to reinforce positive narratives through data‑led proof points and case studies.")
    parts.append("")

    # === 8. Output Mitigate section ===
    parts.append("**Areas Where Additional Context May Help**")
    if mitigate:
        for bucket in sorted(mitigate.keys(), key=sort_key):
            parts.append(mitigate[bucket])
    else:
        parts.append("• Certain narratives may benefit from additional context, transparency or proactive clarification.")
    parts.append("")

    return "\n".join(parts)




# ---------------------------------------------------------
# BUILD STORYBOARD SLIDES
# ---------------------------------------------------------
def compute_overall_score(df_sent, sector_avg: float = 50) -> float:
    """
    Compute an overall sentiment score scaled 0–100 using fair, newspaper-adjusted weights.
    Neutral is excluded from the calculation entirely.
    """

    # Fairer, newspaper-adjusted weights
    weights = {
        "Very Positive": 3.0,
        "Positive": 1.5,
        "Neutral": 0.0,       # excluded from influence
        "Negative": -0.3,
        "Very Negative": -1.5,
    }

    df = df_sent.copy()

    # Map weights
    df["weight_val"] = df["sentiment_display"].map(weights).fillna(0.0)

    # Remove neutral rows entirely from the calculation
    df = df[df["weight_val"] != 0]

    # If everything was neutral or empty → return neutral score
    if df.empty:
        return 50.0

    # Weighted by topic_score if available
    df["weighted"] = df["weight_val"] * df.get("topic_score", 1.0)

    raw = df["weighted"].mean()

    # Scale raw score to 0–100
    max_abs = max(abs(v) for v in weights.values() if v != 0) or 1.0
    scaled = 50.0 + (raw * (50.0 / max_abs))

    return float(max(0.0, min(100.0, scaled)))



def build_overall_gauge_chart(df_sent: pd.DataFrame) -> str:
    """
    Professional semicircular sentiment gauge (0–100).
    50 = neutral (yellow).
    """

    # ------------------------------------------------------------------
    # Compute score safely
    # ------------------------------------------------------------------
    try:
        score = compute_overall_score(df_sent)
    except Exception:
        score = 50.0

    score = max(0.0, min(100.0, float(score)))

    # ------------------------------------------------------------------
    # Figure setup
    # ------------------------------------------------------------------
    fig, ax = plt.subplots(
        figsize=(14, 7),
        subplot_kw={"projection": "polar"}
    )

    bg = "#0E1117"
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)

    ax.set_theta_offset(np.pi)
    ax.set_theta_direction(-1)
    ax.set_ylim(0, 1)

    ax.set_xticks([])
    ax.set_yticks([])
    ax.spines.clear()

    # ------------------------------------------------------------------
    # Smooth Gradient Arc (thicker band, more premium)
    # ------------------------------------------------------------------
    cmap = plt.get_cmap("RdYlGn")
    n = 500
    angles = np.linspace(0, np.pi, n)
    width = angles[1] - angles[0]

    for i, angle in enumerate(angles[:-1]):
        norm_val = i / (n - 1)
        ax.bar(
            angle,
            0.85,           # not full radius (creates inner space)
            width=width,
            bottom=0.15,    # push outward for clean inner area
            color=cmap(norm_val),
            edgecolor="none",
            align="edge",
        )

    # ------------------------------------------------------------------
    # Needle
    # ------------------------------------------------------------------
    needle_angle = np.pi * (score / 100.0)

    ax.plot(
        [needle_angle, needle_angle],
        [0, 0.85],
        color="white",
        linewidth=3,
        zorder=5,
    )

    ax.scatter(
        [needle_angle],
        [0],
        color="white",
        s=90,
        zorder=6,
    )

    # ------------------------------------------------------------------
    # Score Text (closer + cleaner)
    # ------------------------------------------------------------------
    ax.text(
        0.5,
        0.28,  # 🔥 closer to gauge
        f"{score:.1f}",
        transform=ax.transAxes,
        ha="center",
        va="center",
        fontsize=48,
        fontweight="bold",
        color="white",
    )

    ax.text(
        0.5,
        0.20,
        "Overall Score",
        transform=ax.transAxes,
        ha="center",
        va="center",
        fontsize=18,
        color="#C9D1D9",
        fontweight="medium",
    )

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    out = CHART_EXPORT_DIR / "overall_gauge.png"
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor(), dpi=200)
    plt.close(fig)

    return str(out)


def build_storyboard_slides(
    df_sent: pd.DataFrame,
    df_article_sent: pd.DataFrame,
    df_polarity: pd.DataFrame,
    bucket_sizes: pd.DataFrame,
) -> List[SlideSpec]:

    slides: List[SlideSpec] = []

    # --- Slide 1: Executive Overview (no chart) ---
    slides.append(SlideSpec(
        title="Executive Overview",
        body_text=(
            "The media environment over the reporting period reflects a balanced but nuanced sentiment landscape. "
            "While neutral coverage remains the dominant tone, several thematic areas exhibit meaningful positive or "
            "negative skew. This deck summarises the distribution of sentiment, the thematic drivers shaping tone, "
            "and the structural alignment of topics across the corpus."
        ),
        image_path=None,
    ))

    # --- Slide 2: Sentence Sentiment Distribution ---
    chart1 = build_sentence_distribution_chart(df_sent)
    # caption: show percentage by sentiment, highlight dominant tone
    if chart1 is not None:
        chart1 = chart1.properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Sentence-level sentiment distribution — count of sentences by sentiment")
        img1 = save_chart(chart1, "sentence_distribution.png")
        s_counts = df_sent["sentiment_display"].value_counts(normalize=True).mul(100).round(1)
        dominant = s_counts.idxmax()
        dom_pct = s_counts.max()
        caption1 = f"{dominant} is dominant at {dom_pct:.1f}% of sentences; breakdown — " + ", ".join([f"{idx}: {val:.1f}%" for idx, val in s_counts.items()])
    else:
        img1 = None
        caption1 = None
    slides.append(SlideSpec(
        title="Sentence-Level Sentiment Distribution",
        body_text=(
            "Sentence-level sentiment indicates the underlying tonal texture of the corpus. Neutral statements form "
            "the majority, suggesting factual or descriptive reporting. Positive sentiment is present but concentrated "
            "in specific thematic areas, while negative sentiment appears in more targeted clusters."
        ),
        image_path=img1,
        caption=caption1,
    ))

    # --- Slide 3: Article Tone Distribution ---
    chart2 = build_article_tone_chart(df_article_sent)
    if chart2 is not None:
        chart2 = chart2.properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Article-level tone distribution — count of articles by tone")
        img2 = save_chart(chart2, "article_tone_distribution.png")
        if not df_article_sent.empty:
            a_counts = df_article_sent["overall_tone"].value_counts(normalize=True).mul(100).round(1)
            caption2 = f"Article-level distribution — " + ", ".join([f"{idx}: {val:.1f}%" for idx, val in a_counts.items()])
        else:
            caption2 = None
    else:
        img2 = None
        caption2 = None
    slides.append(SlideSpec(
        title="Article-Level Tone Distribution",
        body_text=(
            "When aggregated at the article level, sentiment consolidates into clearer tonal categories. Articles "
            "classified as positive typically contain a higher density of favourable statements, whereas negative "
            "articles reflect sustained critical commentary. The relative balance highlights the structural tone of "
            "media coverage rather than isolated sentence-level fluctuations."
        ),
        image_path=img2,
        caption=caption2,
    ))
    # Ensure bucket labels are wrapped for visibility on charts
    try:
        if df_polarity is not None and not df_polarity.empty:
            def wrap_label(s, n=25):
                if not isinstance(s, str):
                    return s
                if len(s) <= n:
                    return s
                # insert a newline at n and return two-line label
                return s[:n] + "\n" + s[n: n*2]

            df_polarity = df_polarity.copy()
            df_polarity['bucket_short_wrapped'] = df_polarity['bucket_short'].apply(lambda x: wrap_label(x, n=25))
        else:
            df_polarity = df_polarity
    except Exception:
        pass

    # --- Slide 4: Topic Bucket Polarity ---
    chart3 = (
        alt.Chart(df_polarity)
        .mark_bar()
        .encode(
            x=alt.X("bucket_short_wrapped:N", sort="-y", axis=alt.Axis(title=None)),
            y=alt.Y("polarity:Q"),
            color=alt.condition(alt.datum.polarity > 0,
                                alt.value("#2ca02c"),
                                alt.value("#d62728")),
        )
    ).properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Bucket polarity — polarity score by topic bucket")
    img3 = save_chart(chart3, "bucket_polarity.png")
    # caption: top positive and negative buckets + article-level distribution
    article_polarity = compute_article_bucket_polarity(df_sent)
    if not df_polarity.empty:
        top_pos = df_polarity.sort_values("polarity", ascending=False).iloc[0]
        top_neg = df_polarity.sort_values("polarity").iloc[0]
        caption_lines = [f"Top positive bucket: {top_pos['topic_bucket']} ({top_pos['polarity']:.1f}). Top negative bucket: {top_neg['topic_bucket']} ({top_neg['polarity']:.1f})."]
        # If article-level stats available, append percent positive/negative articles for top buckets
        if article_polarity is not None and not article_polarity.empty:
            def art_stats(bucket_name):
                r = article_polarity[article_polarity['topic_bucket'] == bucket_name]
                if r.empty:
                    return None
                row = r.iloc[0]
                return f"Articles — +{row['article_positive_percent']:.0f}% / -{row['article_negative_percent']:.0f}% / ~{row['article_neutral_percent']:.0f}% ({int(row.get('total_articles',0))} articles)"

            pos_stats = art_stats(top_pos['topic_bucket'])
            neg_stats = art_stats(top_neg['topic_bucket'])
            if pos_stats:
                caption_lines.append(f"{top_pos['topic_bucket']} ({pos_stats})")
            if neg_stats:
                caption_lines.append(f"{top_neg['topic_bucket']} ({neg_stats})")

        caption3 = " \n".join(caption_lines)
    else:
        caption3 = None

    slides.append(SlideSpec(
        title="Topic Bucket Polarity",
        body_text=(
            "Polarity scores reveal the directional tone within each topic bucket (sentence-level) — "
            "the chart shows net polarity by bucket. The article-level view (below) summarises the percent of "
            "articles in each bucket that are positive / negative / neutral, which helps surface whether sentence-level "
            "skews are driven by many articles or concentrated sentence mentions. Use combined view to prioritise areas "
            "for communications action."
        ),
        image_path=img3,
        caption=caption3,
    ))

    # --- Slide 5: Coverage Distribution by Topic Bucket ---
    # Short-topic abbreviations for long labels (top long names example)
    short_topic = {
        "Corporate Reputation & Public Perception": "Corp Rep & Perception",
        "Financial Performance & Market Position": "Fin Perf & Market",
        "Customer Experience & Service Delivery": "Cust Exp & Service",
        "Strategy & Transformation": "Strategy & Transform",
        "Workforce, Culture & Operations": "Workforce, Culture & Ops",
        "Products & Offerings": "Products & Offerings",
        "Leadership": "Leadership",
        "Regulation & Compliance": "Regulation & Compliance",
        # placeholders for other long names you may have in top-10
        "Other Long Topic Name 1": "Short Topic 1",
        "Other Long Topic Name 2": "Short Topic 2",
    }

    # Use a display copy so original bucket_sizes remains unchanged
    bucket_sizes_disp = bucket_sizes.copy()
    bucket_sizes_disp["short_topic"] = bucket_sizes_disp["topic_bucket"].map(short_topic).fillna(bucket_sizes_disp["topic_bucket"])

    chart4 = (
        alt.Chart(bucket_sizes_disp)
        .mark_bar()
        .encode(
            x=alt.X(
                "size:Q",
                title="Number of sentences",
            ),
            y=alt.Y(
                "short_topic:N",
                sort="-x",
                axis=alt.Axis(title=None),
            ),
            color=alt.value("#66bd63"),
         tooltip=["topic_bucket", "size"],
        )
    ).properties(
        width=EXPORT_CHART_WIDTH,
        height=EXPORT_CHART_HEIGHT,
        title="Coverage distribution by topic bucket — sentence counts by bucket"
    )

    img4 = save_chart(chart4, "bucket_sizes.png")
    # caption: largest coverage bucket
    if bucket_sizes is not None and not bucket_sizes.empty:
        total = bucket_sizes['size'].sum()
        top = bucket_sizes.sort_values('size', ascending=False).iloc[0]
        caption4 = f"Largest bucket: {top['topic_bucket']} with {int(top['size'])} sentences ({top['size']/total*100:.1f}% of coverage)."
    else:
        caption4 = None
    slides.append(SlideSpec(
        title="Coverage Distribution by Topic Bucket",
        body_text=(
            "Coverage volume varies significantly across buckets, indicating differing levels of media attention. "
            "Larger buckets represent areas of sustained thematic focus, while smaller buckets reflect more episodic "
            "or specialised reporting. Understanding coverage distribution is essential for contextualising sentiment."
        ),
        image_path=img4,
        caption=caption4,
    ))

    # --- Slide 6: Sentiment Composition by Bucket (Heatmap) ---
    heatmap = bucket_sentiment_heatmap(df_sent)
    if heatmap is not None:
        heatmap = heatmap.properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Sentiment composition by bucket — percent share")
        img5 = save_chart(heatmap, "bucket_sentiment_heatmap.png")
        # derive caption from polarity dataframe: highest negative share and highest positive share
        if not df_polarity.empty:
            most_negative = df_polarity.sort_values('negative_percent', ascending=False).iloc[0]
            most_positive = df_polarity.sort_values('positive_percent', ascending=False).iloc[0]
            caption5 = f"Highest negative share: {most_negative['topic_bucket']} ({most_negative['negative_percent']:.1f}%). Highest positive share: {most_positive['topic_bucket']} ({most_positive['positive_percent']:.1f}%)."
        else:
            caption5 = None
    else:
        img5 = None
        caption5 = None
    slides.append(SlideSpec(
        title="Sentiment Composition by Bucket",
        body_text=(
            "The heatmap illustrates the proportional mix of sentiment within each bucket. Buckets with concentrated "
            "negative sentiment may signal emerging risks or persistent critical narratives, while those with strong "
            "positive representation often align with favourable performance, customer experience, or leadership themes."
        ),
        image_path=img5,
        caption=caption5,
    ))

    # --- Slide 7: Sentiment Bubble Chart ---
    bubble = bucket_sentiment_bubble(df_sent)
    if bubble is not None:
        bubble = bubble.properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Sentiment bubble view — share and tonal skew by bucket")
        img6 = save_chart(bubble, "bucket_sentiment_bubble.png")
        # caption: largest bubble (highest percent) and its sentiment
        pivot = (df_sent[df_sent['topic_bucket'].ne('None')].groupby(['topic_bucket','sentiment_display']).size().reset_index(name='count'))
        pivot_tot = pivot.groupby('topic_bucket')['count'].sum().reset_index(name='total')
        pivot = pivot.merge(pivot_tot, on='topic_bucket')
        pivot['percent'] = pivot['count'] / pivot['total'] * 100
        max_row = pivot.loc[pivot['percent'].idxmax()]
        caption6 = f"Largest sentiment share: {max_row['topic_bucket']} — {max_row['sentiment_display']} at {max_row['percent']:.1f}% of its bucket mentions."
    else:
        img6 = None
        caption6 = None
    slides.append(SlideSpec(
        title="Sentiment Bubble View",
        body_text=(
            "Bubble size reflects the share of sentiment within each bucket, while position and colour indicate tonal "
            "skew. Larger bubbles with negative skew warrant closer examination, as they represent high-volume areas "
            "with unfavourable tone. Positive-skewed bubbles highlight strengths and opportunities to reinforce "
            "favourable narratives."
        ),
        image_path=img6,
        caption=caption6,
    ))

    # --- Slide 8: Net Balance vs Intensity ---
    balance = bucket_balance_bubble(df_sent)
    if balance is not None:
        balance = balance.properties(width=EXPORT_CHART_WIDTH, height=EXPORT_CHART_HEIGHT, title="Net balance vs intensity — reputational quadrant view")
        img7 = save_chart(balance, "bucket_balance_bubble.png")
        bb = build_bucket_balance_table(df_sent)
        if not bb.empty:
            top_net = bb.sort_values('net_balance', ascending=False).iloc[0]
            low_net = bb.sort_values('net_balance').iloc[0]
            caption7 = f"Highest net balance: {top_net['topic_bucket']} (net {top_net['net_balance']}, intensity {top_net['avg_intensity']:.2f}). Lowest net balance: {low_net['topic_bucket']} (net {low_net['net_balance']})."
        else:
            caption7 = None
    else:
        img7 = None
        caption7 = None
    slides.append(SlideSpec(
        title="Net Balance vs Sentiment Intensity",
        body_text=(
            "This quadrant view plots buckets by net sentiment balance and average intensity. Buckets in the upper-right "
            "quadrant demonstrate both high positivity and strong sentiment intensity, indicating areas of reputational "
            "strength. Buckets in the lower-left quadrant reflect concentrated negative sentiment and may represent "
            "strategic risk zones requiring targeted intervention."
        ),
        image_path=img7,
        caption=caption7,
    ))

    # --- Slide 9: Topic Drift & Purity ---
    # --- Slide 9: Word Cloud ---
    if WORDCLOUD_EXPORT_PATH.exists():
        # try to compute top term
        try:
            # wordcloud saved from frequencies earlier; attempt to read by inspecting file name not frequencies
            caption_wc = "Word cloud shows user-selected terms sized by mentions; colour encodes sentiment intensity."
        except Exception:
            caption_wc = None
        slides.append(SlideSpec(
            title="Sentiment Word Cloud",
            body_text=(
                "This word cloud visualises the frequency and sentiment of user-selected terms. "
                "Word size reflects mention volume, while colour reflects sentiment and topic-score intensity."
            ),
            image_path=str(WORDCLOUD_EXPORT_PATH),
            caption=caption_wc,
    ))

    # --- Slide 10: Overall Score (gauge) ---
    overall_img = build_overall_gauge_chart(df_sent)
    overall = compute_overall_score(df_sent)
    # interpretative caption
    if overall >= 65:
        interp = "Overall sentiment is strongly positive — favourable coverage dominates."
    elif overall >= 50:
        interp = "Overall sentiment is mildly positive — modest favourable tilt."
    elif overall >= 45:
        interp = "Overall sentiment is neutral — balanced positive and negative signals."
    else:
        interp = "Overall sentiment is negative — unfavourable coverage exceeds favourable."

    # Sector benchmark (placeholder) — allow manual override by editing value
    sector_avg = 50.0  # UK insurers benchmark (placeholder)
    delta = overall - sector_avg
    if abs(delta) < 1.0:
        sector_label = "Neutral"
        sector_note = "vs sector avg"
    elif delta >= 1.0:
        sector_label = "Premium"
        sector_note = f"+{delta:.1f} above avg"
    else:
        sector_label = "Negative"
        sector_note = f"{delta:.1f} below avg"

    caption_text = f"Score: {overall:.1f}/100 ({sector_label} — {sector_note} vs sector avg {sector_avg:.1f}) — {interp}"

    slides.append(SlideSpec(
        title="Overall Sentiment Score",
        body_text=f"The overall sentiment score for this period is {overall:.1f} / 100.",
        image_path=overall_img,
        caption=caption_text,
    ))


    # --- Slide 10: Global Sentiment Drivers ---
    # build driver narrative from full sentence set (includes bucket summaries + top drivers)
    driver_text = build_driver_narrative(df_sent, top_n=3)
    slides.append(SlideSpec(
        title="Global Sentiment Drivers",
        body_text=driver_text,
        image_path=None,  # optional: add a chart if you build one
    ))

    # --- Final Slide: Media & Reputation Action Plan ---
    try:
        overall = compute_overall_score(df_sent)
    except Exception:
        overall = None
    action_plan_text = build_media_reputation_action_plan(df_polarity, overall, df_sent)
    slides.append(SlideSpec(
        title="Media & Reputation Action Plan",
        body_text=action_plan_text,
        image_path=None,
    ))

    return slides

# ---------------------------------------------------------
# EXPORT STORYBOARD TO PPTX
# ---------------------------------------------------------

# --- DESIGN CONSTANTS (TUNE FREELY) ---
TITLE_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(13)
TOP_MARGIN = Inches(1.0)
SIDE_MARGIN = Inches(0.7)
IMAGE_TOP = Inches(2.0)
IMAGE_HEIGHT = Inches(5.0)  # enlarged for PPTX slides
LINE_SPACING = 1.2  # multiplier


def export_storyboard_to_pptx(slides: List[SlideSpec]) -> bytes:
    prs = Presentation()

    for slide_spec in slides:
        layout = prs.slide_layouts[5]  # Title Only layout (usually index 5)
        slide = prs.slides.add_slide(layout)
        fill = slide.background.fill 
        fill.solid()
        fill.fore_color.rgb = RGBColor(14, 17, 23)
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_spec.title
        title_tf = title_shape.text_frame
        title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
        title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        # Body text placeholder (manual textbox for full control)
        left = SIDE_MARGIN
        top = TOP_MARGIN + IMAGE_HEIGHT + Inches(0.4)
        width = prs.slide_width - 2 * SIDE_MARGIN
        height = prs.slide_height - top - Inches(0.5)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        body_tf = body_box.text_frame
        body_tf.clear()
        p = body_tf.paragraphs[0]
        p.text = slide_spec.body_text
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        body_tf.word_wrap = True
        body_tf.auto_size = False

        # Optional image
        if slide_spec.image_path is not None and os.path.exists(slide_spec.image_path):
            # Size the chart image for PPTX and preserve aspect ratio by specifying height only
            chart_height = Inches(5.5)
            img_left = SIDE_MARGIN
            img_top = TOP_MARGIN

            pic = slide.shapes.add_picture(
                slide_spec.image_path,
                img_left,
                img_top,
                height=chart_height,
            )
            # pic.width/height now reflect the placed image size; center caption using available width
            img_width = pic.width

            # Optional caption under image (short, data-specific)
            if getattr(slide_spec, 'caption', None):
                cap_top = img_top + chart_height + Inches(0.1)
                cap_box = slide.shapes.add_textbox(SIDE_MARGIN, cap_top, img_width, Inches(0.4))
                cap_tf = cap_box.text_frame
                cap_tf.clear()
                cap_p = cap_tf.paragraphs[0]
                cap_p.text = slide_spec.caption
                cap_p.font.size = Pt(13)
                cap_p.font.italic = True
                cap_p.font.color.rgb = RGBColor(180, 180, 180)
                cap_p.alignment = PP_ALIGN.LEFT
    

# Set slide background to dark blue
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(14, 17, 23)  # same as Streamlit dark mode (#0E1117)
     # same as Streamlit dark mode (#0E1117)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def render_powerpoint_storyboard(df_sent, df_topics, bucket_sizes):
    st.header("PowerPoint Storyboard")

    df_article_sent = compute_article_sentiment(df_sent)
    df_polarity = compute_bucket_polarity(df_sent)

    slides = build_storyboard_slides(
        df_sent=df_sent,
        df_article_sent=df_article_sent,
        df_polarity=df_polarity,
        bucket_sizes=bucket_sizes,
    )

    st.markdown("### Download Storyboard as PowerPoint")
    pptx_bytes = export_storyboard_to_pptx(slides)
    st.download_button(
        label="Download PPTX",
        data=pptx_bytes,
        file_name="media_intelligence_storyboard.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

def render_narrative_export(df_sent):
    st.header("Narrative Export")
    narrative = build_narrative(df_sent)
    st.write(narrative)
    st.download_button(
        label="Download Narrative",
        data=narrative,
        file_name="media_intelligence_narrative.txt",
        mime="text/plain",
    )

# ---------------------------------------------------------
# MAIN APP
# ---------------------------------------------------------

def main():
    st.set_page_config(
        page_title="Media Intelligence Dashboard",
        layout="wide",
    )

    st.title("Media Intelligence Dashboard")

    (
        data,
        df_sent,
        df_topics,
        df_entities,
        df_linked,
        df_ent_sent,
        df_ent_time,
    ) = load_master()

    df_sent, df_topics, bucket_sizes = apply_topic_buckets(df_sent, df_topics)
    df_articles = pd.DataFrame(data.get("articles", []))

    # Diagnostics
    st.sidebar.write("Sentences columns:", df_sent.columns.tolist())
    st.sidebar.write("Topics columns:", df_topics.columns.tolist())

    page = st.sidebar.selectbox(
        "Page",
        [
            "Executive Summary",
            "Topic Buckets",
            "Overview",
            "Topic Explorer",
            "Sentence Inspector",
            "Entity Explorer",
            "Article Browser",
            "Search",
            "PowerPoint Storyboard",
            "Narrative Export",
        ],
    )

    if page == "Executive Summary":
        render_executive_summary_page(df_sent)

    elif page == "Topic Buckets":
        render_topic_buckets_page(df_sent, df_topics, bucket_sizes)

    elif page == "Overview":
        render_overview_page(df_sent, df_topics, df_articles)

    elif page == "Topic Explorer":
        render_topic_explorer_page(df_sent)

    elif page == "Sentence Inspector":
        render_sentence_inspector_page(df_sent)

    elif page == "Entity Explorer":
        render_entity_explorer_page(df_entities, df_sent)

    elif page == "Article Browser":
        render_article_browser_page(df_articles, df_sent)

    elif page == "Search":
        render_search_page(df_sent, df_articles)
    
    elif page == "PowerPoint Storyboard":
        render_powerpoint_storyboard(df_sent, df_topics, bucket_sizes)
    
    elif page == "Narrative Export":
        render_narrative_export(df_sent)


if __name__ == "__main__":
    main()
