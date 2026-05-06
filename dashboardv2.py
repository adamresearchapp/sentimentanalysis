# ---------------------------------------------------------
# Media Intelligence Dashboard (dashboardV2.py)
# Hardened version aligned with databuilder.py
# Design System Standardisation + Analytics Clarity + Topic Classification & Governance
# ---------------------------------------------------------

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
try:
    from pptx.oxml import OxmlElement
except ImportError:
    from pptx.oxml.xmlchemy import OxmlElement
import io
import hmac
from pathlib import Path
from dataclasses import dataclass
from typing import List, Optional
import altair as alt
import pandas as pd
import streamlit as st
from wordcloud import WordCloud 
import matplotlib.pyplot as plt
import numpy as np
import os
import shutil
import time
import re
import sys

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import vl_convert as vlc
except Exception:
    vlc = None

# ---------------------------------------------------------
# DESIGN SYSTEM STANDARDISATION (Point 1)
# ---------------------------------------------------------
# Consistent typography: Helvetica Neue for titles, Garamond for body
# Unified colour palette: transparent chart canvases with configurable blue accents
# Cleaner visual hierarchy: Standardized axis, grid, and sizes
# Optimised sizing: EXPORT_CHART_WIDTH/HEIGHT tuned for slides

PRIMARY_BLUE = "#2A3653"
SECONDARY_BLUE = "#48708A"
GRID_COLOR = "#E0E0E0"

# Altair view / export PNG canvas (Streamlit page or slide background shows through).
CHART_VIEW_FILL = "transparent"

PPT_PRIMARY_RGB = RGBColor(42, 54, 83)
PPT_SECONDARY_RGB = RGBColor(72, 112, 138)

def axa_white_refined_theme():
    """
    Professional theme with transparent chart views, Helvetica/Garamond fonts,
    and consistent blue accents for axes, legends, and titles.
    """
    return {
        "config": {
            "background": CHART_VIEW_FILL,
            "view": {
                "fill": CHART_VIEW_FILL,
                "stroke": "transparent",
                "continuousWidth": EXPORT_CHART_WIDTH,
                "continuousHeight": EXPORT_CHART_HEIGHT,
            },
            "axis": {
                "labelFont": "Garamond",
                "labelFontSize": 18,
                "titleFont": "Helvetica Neue",
                "titleFontSize": 22,
                "labelAngle": 0,
                "labelAlign": "center",
                "labelBaseline": "middle",
                "labelPadding": 12,
                "labelOverlap": "greedy",
                "labelLimit": 1000,
                "titleLimit": 1000,
                "titlePadding": 10,
                "labelColor": SECONDARY_BLUE,
                "titleColor": PRIMARY_BLUE,
                "gridColor": GRID_COLOR,
            },
            "legend": {
                "labelFont": "Garamond",
                "labelFontSize": 18,
                "titleFont": "Helvetica Neue",
                "titleFontSize": 22,
                "labelColor": SECONDARY_BLUE,
                "titleColor": PRIMARY_BLUE,
                "labelLimit": 1000,
                "titleLimit": 1000,
            },
            "title": {
                "font": "Helvetica Neue",
                "fontSize": 28,
                "color": PRIMARY_BLUE,
                "anchor": "start",
                "offset": 10,
            },
            "text": {
                "font": "Garamond",
                "fontSize": 13,
                "color": SECONDARY_BLUE,
            },
        }
    }

# Register and enable the AXA white refined theme
alt.themes.register("axa_white_refined", axa_white_refined_theme)
alt.themes.enable("axa_white_refined")

# ---------------------------------------------------------
# CHART EXPORT DIRECTORY & SIZING CONSTANTS
# ---------------------------------------------------------

CHART_EXPORT_DIR = Path("./powerpoint")
CHART_EXPORT_DIR.mkdir(parents=True, exist_ok=True)

DEFAULT_CHART_WIDTH = 720
DEFAULT_CHART_HEIGHT = 420
EXPORT_CHART_WIDTH = 1200
EXPORT_CHART_HEIGHT = 700
EXPORT_SCALE_FACTOR = 2

AXIS_LABEL_PADDING = 18
AXIS_TITLE_PADDING = 24
AXIS_TITLE_LIMIT = 1000

WORDCLOUD_EXPORT_PATH = CHART_EXPORT_DIR / "wordcloud.png"

# ---------------------------------------------------------
# PATH TO MASTER JSON (must match sentimentanalysis.py default)
# ---------------------------------------------------------

_PROJECT_ROOT = Path(__file__).resolve().parent
MASTER_JSON = Path(os.getenv("PRESSCHOICE_MASTER_JSON", str(_PROJECT_ROOT / "master.json")))

# ---------------------------------------------------------
# SENTIMENT LABEL CLEANUP
# ---------------------------------------------------------

SENTIMENT_LABEL_DISPLAY = {
    "very_negative": "Very Negative",
    "negative": "Negative",
    "neutral": "Neutral",
    "positive": "Positive",
    "very_positive": "Very Positive",
}

SENTIMENT_ORDER = [
    "Very Negative",
    "Negative",
    "Neutral",
    "Positive",
    "Very Positive",
]

SENTIMENT_WEIGHTS = {
    "Very Positive": 2,
    "Positive": 1,
    "Neutral": 0,
    "Negative": -1,
    "Very Negative": -2,
}

# Speedometer-only calibration. Financial journalism often has a cautious/negative
# baseline, so the executive gauge gives positive signals a modest uplift without
# altering raw labels, distribution charts, or bucket polarity.
GAUGE_SENTIMENT_WEIGHTS = {
    "Very Positive": 2.25,
    "Positive": 1.25,
    "Neutral": 0,
    "Negative": -1,
    "Very Negative": -2,
}

SENTIMENT_COLORS = {
    "Very Negative": "#d73027",
    "Negative": "#fc8d59",
    "Neutral": "#808080",
    "Positive": "#66bd63",
    "Very Positive": "#1a9850",
}

# ---------------------------------------------------------
# TOPIC DEFINITIONS & GOVERNANCE FRAMEWORK (Point 3)
# ---------------------------------------------------------

try:
    from sentimentanalysis import (
        LOW_SENTIMENT_CONFIDENCE,
        LOW_TOPIC_CONFIDENCE,
        TOPIC_DEFINITIONS,
        TOPIC_MARGIN_DRIFT,
        TOPIC_NEAR_TIE_MARGIN,
        update_sentence_override,
    )
except Exception:
    def update_sentence_override(master: dict, global_index: int, sentiment_override: str = None, topic_override: str = None) -> bool:
        """Fallback: no-op for environments where sentimentanalysis isn't importable."""
        for sentence in master.get("sentences", []):
            if sentence.get("global_index") == global_index:
                if sentiment_override:
                    sentence["manual_sentiment_override"] = sentiment_override
                if topic_override:
                    sentence["manual_topic_override"] = topic_override
                return True
        return False
    
    TOPIC_DEFINITIONS = {
        "Leadership & Governance": ["CEO", "executive", "leadership", "governance"],
        "Financial Performance & Market Position": ["earnings", "revenue", "profit", "market"],
        "Customer Experience & Service Delivery": ["customer", "service", "experience"],
        "Products & Offerings": ["product", "offering", "launch"],
        "Corporate Reputation & Public Perception": ["reputation", "brand", "perception"],
        "Strategy & Transformation": ["strategy", "transformation", "innovation"],
        "Regulation & Compliance": ["regulation", "compliance", "regulatory"],
        "Workforce, Culture & Operations": ["workforce", "culture", "operations"],
    }
    LOW_TOPIC_CONFIDENCE = 0.25
    LOW_SENTIMENT_CONFIDENCE = 0.30
    TOPIC_NEAR_TIE_MARGIN = 0.08
    TOPIC_MARGIN_DRIFT = 0.10

# ---------------------------------------------------------
# TOPIC BUCKET MAPPING
# ---------------------------------------------------------

TOPIC_BUCKET_MAP = {
    # Performance & Strategy
    "Financial Performance & Market Position": "Performance & Strategy",
    "Strategy & Transformation": "Performance & Strategy",

    # Customer & Brand Experience
    "Customer Experience & Service Delivery": "Customer & Brand Experience",
    "Corporate Reputation & Public Perception": "Customer & Brand Experience",
    "Products & Offerings": "Customer & Brand Experience",

    # Governance, Leadership & Accountability
    "Leadership & Governance": "Governance, Leadership & Accountability",
    "Leadership": "Governance, Leadership & Accountability",
    "Regulation & Compliance": "Governance, Leadership & Accountability",

    # Workforce, Culture & Operations
    "Workforce, Culture & Operations": "Workforce, Culture & Operations",
}

BUCKET_SHORT = {
    "Performance & Strategy": "PS",
    "Customer & Brand Experience": "CB",
    "Governance, Leadership & Accountability": "GL",
    "Workforce, Culture & Operations": "WO",
}

# Short codes for chart axes only (full names remain in tooltips / tables).
TOPIC_GRAPH_LABEL = {
    "Corporate Reputation & Public Perception": "CR",
    "Customer Experience & Service Delivery": "CX",
    "Financial Performance & Market Position": "FP",
    "Leadership": "LG",
    "Leadership & Governance": "LG",
    "Products & Offerings": "PO",
    "Regulation & Compliance": "RC",
    "Strategy & Transformation": "ST",
    "Workforce, Culture & Operations": "WF",
}

BUCKET_ORDER = [
    "Customer & Brand Experience",
    "Governance, Leadership & Accountability",
    "Performance & Strategy",
    "Workforce, Culture & Operations",
]

# Inverse of BUCKET_SHORT for charts / tooltips.
BUCKET_SHORT_TO_NAME = {abbrev: full for full, abbrev in BUCKET_SHORT.items()}

TOPIC_THRESHOLD = 0.18
TOPIC_ALPHA_EMBED = 0.6
TOPIC_ALPHA_NLI = 0.4

BUCKET_KEY_TEXT = (
    "Buckets key: "
    "CB = Customer & Brand Experience, "
    "GL = Governance, Leadership & Accountability, "
    "PS = Performance & Strategy, "
    "WO = Workforce, Culture & Operations."
)

FINE_TOPIC_KEY_TEXT = (
    "Fine topics key: "
    "CR = Corporate Reputation & Public Perception, "
    "CX = Customer Experience & Service Delivery, "
    "FP = Financial Performance & Market Position, "
    "LG = Leadership & Governance, "
    "PO = Products & Offerings, "
    "RC = Regulation & Compliance, "
    "ST = Strategy & Transformation, "
    "WF = Workforce, Culture & Operations."
)

PPT_FOOTER_KEY_TEXT = (
    "Key: Buckets CB=Customer & Brand Experience, GL=Governance/Leadership/Accountability, "
    "PS=Performance & Strategy, WO=Workforce/Culture/Operations | "
    "Fine topics CR=Corporate Reputation, CX=Customer Experience, FP=Financial Performance, "
    "LG=Leadership & Governance, PO=Products & Offerings, RC=Regulation & Compliance, "
    "ST=Strategy & Transformation, WF=Workforce/Culture/Operations"
)

# ---------------------------------------------------------
# DESIGN FUNCTIONS: SAVE CHART WITH THEME
# ---------------------------------------------------------

def save_chart(chart: alt.Chart, filename: str) -> str:
    """
    Export an Altair chart to CHART_EXPORT_DIR as PNG with AXA theme applied.
    Uses a transparent view so the Streamlit / slide background shows through.
    Returns file path string (may not exist if save failed, but caller checks).
    """
    if not filename.lower().endswith(".png"):
        filename = f"{filename}.png"
    out_path = CHART_EXPORT_DIR / filename
    svg_path = out_path.with_suffix(".svg")

    # Avoid stale artifacts: always remove previous exports first.
    try:
        if out_path.exists():
            out_path.unlink()
        if svg_path.exists():
            svg_path.unlink()
    except Exception:
        pass

    try:
        # Keep chart content identical to the dashboard, but make export bounds
        # include titles/axis labels so PowerPoint PNGs are self-contained.
        chart = chart.properties(
            width=EXPORT_CHART_WIDTH,
            height=EXPORT_CHART_HEIGHT,
        ).configure_axis(
            titleLimit=AXIS_TITLE_LIMIT,
            titlePadding=AXIS_TITLE_PADDING,
            labelPadding=AXIS_LABEL_PADDING,
            labelOverlap=False,
            labelColor="#111111",
            titleColor="#111111",
            labelFont="sans-serif",
            titleFont="sans-serif",
            labelFontSize=14,
            titleFontSize=16,
        ).configure_legend(
            titleLimit=AXIS_TITLE_LIMIT,
            labelLimit=AXIS_TITLE_LIMIT,
            labelColor="#111111",
            titleColor="#111111",
            labelFont="sans-serif",
            titleFont="sans-serif",
            labelFontSize=12,
            titleFontSize=14,
        ).configure_title(
            color="#111111",
            font="sans-serif",
            fontSize=20,
        ).configure_view(
            stroke="#DDDDDD",
            fill="white",
        ).configure(
            background="white",
        ).properties(
            autosize=alt.AutoSizeParams(type="pad", contains="padding"),
        )

        # Preferred path on hosted environments: render via vl-convert directly.
        if vlc is not None:
            spec = chart.to_dict()
            png_bytes = vlc.vegalite_to_png(spec, scale=EXPORT_SCALE_FACTOR)
            out_path.write_bytes(png_bytes)
            return str(out_path)

        # Fallback when vl-convert is unavailable.
        chart.save(str(out_path), scale_factor=EXPORT_SCALE_FACTOR)
        return str(out_path)
    except Exception:
        # Fallback: attempt SVG then convert if possible, else save basic PNG
        try:
            chart.save(str(svg_path))
        except Exception:
            try:
                chart.save(str(out_path))
            except Exception:
                # Last resort: create a minimal placeholder PNG
                try:
                    fig = plt.figure(figsize=(EXPORT_CHART_WIDTH / 100, EXPORT_CHART_HEIGHT / 100), dpi=100)
                    fig.patch.set_facecolor("none")
                    plt.text(
                        0.5, 0.5,
                        "Chart export failed",
                        ha="center",
                        va="center",
                        fontsize=14,
                        color="#666666",
                    )
                    plt.axis("off")
                    fig.savefig(out_path, bbox_inches="tight", facecolor="none", transparent=True, dpi=150)
                    plt.close(fig)
                except Exception:
                    pass
        return str(out_path)


def _save_mpl_export(fig, filename: str) -> str:
    if not filename.lower().endswith(".png"):
        filename = f"{filename}.png"
    out_path = CHART_EXPORT_DIR / filename
    try:
        if out_path.exists():
            out_path.unlink()
    except Exception:
        pass
    fig.patch.set_facecolor("none")
    fig.patch.set_alpha(0)
    for ax in fig.axes:
        ax.set_facecolor("none")
        ax.patch.set_alpha(0)
    fig.tight_layout(pad=0.6)
    fig.savefig(out_path, dpi=180, bbox_inches="tight", facecolor="none", edgecolor="none", transparent=True)
    plt.close(fig)
    return str(out_path)


def _style_mpl_axis(ax, title: str, xlabel: str = "", ylabel: str = ""):
    # Slide titles already describe the chart; omit in-PNG titles to maximise plot area.
    ax.set_title("")
    ax.set_xlabel(xlabel, color=PRIMARY_BLUE, fontsize=11, labelpad=10)
    ax.set_ylabel(ylabel, color=PRIMARY_BLUE, fontsize=11, labelpad=10)
    ax.tick_params(axis="both", colors=SECONDARY_BLUE, labelsize=10)
    ax.grid(axis="y", color=GRID_COLOR, linewidth=0.8, alpha=0.9)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#B8B8B8")
    ax.spines["bottom"].set_color("#B8B8B8")


def export_sentence_distribution_png(df_sent: pd.DataFrame, filename: str) -> str | None:
    if df_sent is None or df_sent.empty:
        return None
    df = df_sent.groupby("sentiment_display").size().reindex(SENTIMENT_ORDER, fill_value=0)
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.bar(df.index, df.values, color=[SENTIMENT_COLORS[s] for s in df.index])
    _style_mpl_axis(ax, "Sentence Sentiment Distribution", "Sentiment", "Number of sentences")
    ax.tick_params(axis="x", rotation=0)
    return _save_mpl_export(fig, filename)


def export_article_tone_png(df_article_sent: pd.DataFrame, filename: str) -> str | None:
    if df_article_sent is None or df_article_sent.empty:
        return None
    order = ["Negative", "Neutral", "Positive"]
    colors = [SENTIMENT_COLORS["Negative"], SENTIMENT_COLORS["Neutral"], SENTIMENT_COLORS["Positive"]]
    df = df_article_sent.groupby("overall_tone").size().reindex(order, fill_value=0)
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.bar(df.index, df.values, color=colors)
    _style_mpl_axis(ax, "Article Tone Distribution", "Article tone", "Number of articles")
    return _save_mpl_export(fig, filename)


def export_bucket_sizes_png(bucket_sizes: pd.DataFrame, filename: str) -> str | None:
    if bucket_sizes is None or bucket_sizes.empty or "topic_bucket" not in bucket_sizes.columns:
        return None
    df = bucket_sizes[bucket_sizes["topic_bucket"].ne("None")].copy()
    if df.empty:
        return None
    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]
    series = df.set_index("bucket_short")["size"].reindex(domain).dropna()
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.bar(series.index, series.values, color=SECONDARY_BLUE)
    _style_mpl_axis(ax, "Bucket Sizes", "Bucket (code)", "Sentences")
    return _save_mpl_export(fig, filename)


def export_bucket_polarity_png(df_polarity: pd.DataFrame, filename: str) -> str | None:
    if df_polarity is None or df_polarity.empty:
        return None
    df = df_polarity.copy()
    if "bucket_short" not in df.columns:
        df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]
    df = df.set_index("bucket_short").reindex(domain).dropna(subset=["polarity"]).reset_index()
    fig, ax = plt.subplots(figsize=(8.6, 5.0))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    colors = ["#2ca02c" if v > 0 else "#d62728" for v in df["polarity"]]
    ax.bar(df["bucket_short"], df["polarity"], color=colors)
    ax.axhline(0, color="#555555", linewidth=1)
    _style_mpl_axis(ax, "Bucket Polarity", "Bucket (code)", "Polarity score")
    return _save_mpl_export(fig, filename)


def export_bucket_balance_png(df_sent: pd.DataFrame, filename: str) -> str | None:
    table = build_bucket_balance_table(df_sent)
    if table is None or table.empty:
        return None
    fig, ax = plt.subplots(figsize=(9.2, 5.2))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    sizes = np.interp(table["total_count"], (table["total_count"].min(), table["total_count"].max()), (200, 1200)) if table["total_count"].nunique() > 1 else np.repeat(650, len(table))
    sc = ax.scatter(table["net_balance"], table["avg_intensity"], s=sizes, c=table["avg_intensity"], cmap="RdYlGn", alpha=0.75, edgecolors="black", linewidths=0.5)
    for _, row in table.iterrows():
        ax.annotate(row["bucket_short"], (row["net_balance"], row["avg_intensity"]), xytext=(8, 3), textcoords="offset points", color=PRIMARY_BLUE, fontsize=11, fontweight="bold")
    ax.axhline(0, color="#555555", linewidth=1.2)
    ax.axvline(0, color="#555555", linewidth=1.2)
    _style_mpl_axis(ax, "Bucket Balance Map", "Net balance", "Intensity")
    return _save_mpl_export(fig, filename)


def _bucket_sentiment_pivot(df_sent: pd.DataFrame) -> pd.DataFrame:
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty:
        return pd.DataFrame()
    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    counts = df.groupby(["bucket_short", "sentiment_display"]).size().reset_index(name="count")
    totals = counts.groupby("bucket_short")["count"].sum().rename("total")
    counts = counts.merge(totals, on="bucket_short")
    counts["percent"] = counts["count"] / counts["total"] * 100
    matrix = counts.pivot(index="bucket_short", columns="sentiment_display", values="percent")
    return matrix.reindex(index=[BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT], columns=SENTIMENT_ORDER).fillna(0)


def export_bucket_sentiment_heatmap_png(df_sent: pd.DataFrame, filename: str) -> str | None:
    matrix = _bucket_sentiment_pivot(df_sent)
    if matrix.empty:
        return None
    fig, ax = plt.subplots(figsize=(8.8, 5.2))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    cmap = plt.colormaps["Blues"].copy()
    cmap.set_bad((1, 1, 1, 0))
    masked_values = np.ma.masked_where(matrix.values == 0, matrix.values)
    im = ax.imshow(masked_values, cmap=cmap, aspect="auto")
    ax.set_xticks(range(len(matrix.columns)), matrix.columns, rotation=0)
    ax.set_yticks(range(len(matrix.index)), matrix.index)
    _style_mpl_axis(ax, "Bucket x Sentiment Composition", "Sentiment", "Bucket (code)")
    fig.colorbar(im, ax=ax, fraction=0.035, pad=0.03, label="Percent of bucket")
    return _save_mpl_export(fig, filename)


def export_bucket_sentiment_bubble_png(df_sent: pd.DataFrame, filename: str) -> str | None:
    matrix = _bucket_sentiment_pivot(df_sent)
    if matrix.empty:
        return None
    rows = []
    for yi, bucket in enumerate(matrix.index):
        for xi, sent in enumerate(matrix.columns):
            rows.append({"x": xi, "y": yi, "bucket": bucket, "sentiment": sent, "percent": matrix.loc[bucket, sent]})
    d = pd.DataFrame(rows)
    fig, ax = plt.subplots(figsize=(9.8, 5.8))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    colors = [SENTIMENT_COLORS[s] for s in d["sentiment"]]
    ax.scatter(d["x"], d["y"], s=d["percent"].clip(lower=0) * 18 + 20, c=colors, alpha=0.85, edgecolors="black", linewidths=0.3)
    ax.set_xticks(range(len(matrix.columns)), matrix.columns)
    ax.set_yticks(range(len(matrix.index)), matrix.index)
    ax.invert_yaxis()
    ax.set_xlim(-0.35, len(matrix.columns) - 0.65)
    ax.set_ylim(len(matrix.index) - 0.55, -0.45)
    ax.tick_params(axis="x", pad=8)
    ax.tick_params(axis="y", pad=8)
    _style_mpl_axis(ax, "Bucket x Sentiment Bubble View", "Sentiment", "Bucket (code)")
    return _save_mpl_export(fig, filename)


def export_topic_drift_heatmap_png(df_sent: pd.DataFrame, filename: str) -> str | None:
    if "topic_bucket" not in df_sent.columns or "topic_name" not in df_sent.columns:
        return None
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty:
        return None
    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    df["topic_graph"] = df["topic_name"].map(TOPIC_GRAPH_LABEL).fillna(df["topic_name"].astype(str).str.slice(0, 24))
    counts = df.groupby(["topic_graph", "bucket_short"]).size().reset_index(name="count")
    totals = df.groupby("bucket_short").size().rename("bucket_total")
    counts = counts.merge(totals, on="bucket_short")
    counts["percent"] = counts["count"] / counts["bucket_total"] * 100
    bucket_domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]
    topic_domain = [TOPIC_GRAPH_LABEL[t] for t in TOPIC_DEFINITIONS.keys() if t in TOPIC_GRAPH_LABEL and TOPIC_GRAPH_LABEL[t] in set(counts["topic_graph"])]
    matrix = counts.pivot(index="topic_graph", columns="bucket_short", values="percent").reindex(index=topic_domain, columns=bucket_domain).fillna(0)
    if matrix.empty:
        return None
    fig, ax = plt.subplots(figsize=(8.8, max(5.0, 0.45 * len(matrix.index) + 1.8)))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    cmap = plt.colormaps["Greens"].copy()
    cmap.set_bad((1, 1, 1, 0))
    masked_values = np.ma.masked_where(matrix.values == 0, matrix.values)
    im = ax.imshow(masked_values, cmap=cmap, aspect="auto", vmin=0, vmax=max(100, float(matrix.values.max())))
    ax.set_xticks(range(len(matrix.columns)), matrix.columns)
    ax.set_yticks(range(len(matrix.index)), matrix.index)
    _style_mpl_axis(ax, "Topic mix within each bucket (% of bucket sentences)", "Bucket (code)", "Fine topic (code)")
    fig.colorbar(im, ax=ax, fraction=0.035, pad=0.03, label="% of bucket")
    return _save_mpl_export(fig, filename)


def export_topic_salience_png(df_sent: pd.DataFrame, filename: str, top_n: int = 8) -> str | None:
    salience = compute_topic_salience(df_sent)
    if salience is None or salience.empty:
        return None
    df = salience.head(top_n).copy()
    df["topic_graph"] = df["topic_name"].map(TOPIC_GRAPH_LABEL).fillna(df["topic_name"])
    fig, ax = plt.subplots(figsize=(8.8, max(5.0, 0.45 * len(df) + 1.8)))
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    ax.barh(df["topic_graph"], df["pct_of_bucket"], color=SECONDARY_BLUE)
    ax.invert_yaxis()
    ax.set_xlim(0, 100)
    _style_mpl_axis(ax, "Fine Topic Salience", "% of bucket sentences (salience)", "Fine topic (code)")
    return _save_mpl_export(fig, filename)



# ---------------------------------------------------------
# LOAD MASTER JSON
# ---------------------------------------------------------

@st.cache_data
def load_master():
    data = json.loads(MASTER_JSON.read_text(encoding="utf-8"))

    df_sent = pd.DataFrame(data.get("sentences", []))
    df_topics = pd.DataFrame(data.get("topics", []))
    df_entities = pd.DataFrame(data.get("entities_corpus", []))
    df_linked = pd.DataFrame(data.get("entities_linked", []))
    df_ent_sent = pd.DataFrame(data.get("entity_sentiment", []))
    df_ent_time = pd.DataFrame(data.get("entity_timeline", []))

    # Fallbacks to avoid KeyError
    if "topic_name" not in df_sent.columns:
        df_sent["topic_name"] = "None"
    if "topic_score" not in df_sent.columns:
        df_sent["topic_score"] = 0.0

    if "topic_name" not in df_topics.columns and not df_topics.empty:
        df_topics["topic_name"] = df_sent["topic_name"]
    if "topic_score" not in df_topics.columns and not df_topics.empty:
        df_topics["topic_score"] = df_sent["topic_score"]

    # Always create sentiment_display
    if "label_5" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["label_5"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_label" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_label"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_class" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_class"].map(SENTIMENT_LABEL_DISPLAY)
    elif "sentiment_category" in df_sent.columns:
        df_sent["sentiment_display"] = df_sent["sentiment_category"].map(SENTIMENT_LABEL_DISPLAY)
    else:
        df_sent["sentiment_display"] = "Neutral"

    return data, df_sent, df_topics, df_entities, df_linked, df_ent_sent, df_ent_time

# ---------------------------------------------------------
# SECTION 6 — UPDATED: Sentence Inspector with Manual Overrides
# ---------------------------------------------------------

def render_sentence_inspector_page(df_sent):
    st.header("Sentence Inspector")

    query = st.text_input("Search text in sentences (single-term)")
    multi_terms = st.text_input("Word Cloud Terms (comma-separated)")

    # --- Manual Override Section (GOVERNANCE FRAMEWORK: Point 3) ---
    st.subheader("Manual Classification Overrides")
    st.write("""
    Use this section to review and manually correct sentence-level classifications.
    Select a sentence by `global_index`, choose new sentiment/topic labels, and apply.
    Changes are saved to `master.json` and the dashboard reloads: charts and exports use the updated rows immediately.
    Run the full pipeline only when you add or remove articles or want to recompute model outputs from scratch.
    """)
    
    render_manual_override_controls(df_sent, key_prefix="inspector")

    # --- Table Search ---
    st.subheader("Search & Review")
    if query:
        df_q = df_sent[
            df_sent["sentence"].str.contains(query, case=False, na=False)
        ].copy()

        st.write(f"Found {len(df_q)} matching sentences.")
        
        # Show review flags if available
        if "needs_review" in df_q.columns:
            needs_review_count = (df_q["needs_review"] == True).sum()
            if needs_review_count > 0:
                st.warning(f"⚠ {needs_review_count} sentences flagged for manual review (low confidence or leadership mention).")
                df_q_review = df_q[df_q["needs_review"] == True]
                st.dataframe(
                    df_q_review[[
                        "global_index",
                        "sentence",
                        "sentiment_display",
                        "topic_name",
                        "topic_bucket",
                        "review_reasons",
                    ]]
                )
        
        st.dataframe(
            df_q[
                [
                    "global_index",
                    "sentence",
                    "topic_name",
                    "topic_bucket",
                    "sentiment_display",
                    "topic_score",
                ]
            ]
        )
    else:
        st.write("Enter text above to search through all sentences.")

    # --- WORD CLOUD ---
    st.subheader("Sentiment Word Cloud")
    if multi_terms:
        freq, colors = build_sentiment_wordcloud_data(df_sent, multi_terms)
        if freq:
            wc = WordCloud(
                width=1200,
                height=600,
                background_color=None,
                mode="RGBA",
                prefer_horizontal=1.0,
            ).generate_from_frequencies(freq)

            wc.recolor(color_func=make_color_func(colors))

            fig, ax = plt.subplots(figsize=(12, 6))
            fig.patch.set_facecolor("none")
            ax.set_facecolor("none")
            ax.imshow(wc, interpolation="bilinear")
            ax.axis("off")

            # Save for PPTX export
            wc.to_file(str(WORDCLOUD_EXPORT_PATH))

            st.pyplot(fig)
        else:
            st.info("No matching sentences for word cloud.")


def _review_reasons_as_list(val) -> list:
    if isinstance(val, list):
        return val
    return []


def _leadership_keyword_only(reasons: list) -> bool:
    return bool(reasons) and set(reasons) == {"leadership_figure"}


def compute_soft_outlier_masks(
    df: pd.DataFrame,
    *,
    mode: str,
    tail_pct: int,
    z_mult: float,
) -> tuple[pd.Series, pd.Series, pd.Series, dict]:
    """
    Soft queue rules (dashboard-only) on topic_confidence, sentiment_confidence, topic_margin.

    Hybrid topic scores are 0–1 after normalising across eight topics, so absolute cutoffs
    like 0.35 flag almost everything. Corpus-relative modes compare each sentence to the
    distribution of the current export instead.
    """
    tc = pd.to_numeric(df["topic_confidence"], errors="coerce")
    sc = pd.to_numeric(df["sentiment_confidence"], errors="coerce")
    tm = pd.to_numeric(df["topic_margin"], errors="coerce")

    meta: dict = {"mode": mode}
    n = int(len(df))

    if mode == "fixed":
        meta.update({"topic_thr": 0.35, "sent_thr": 0.35, "margin_thr": 0.10})
        return (
            tc.fillna(0.0) < 0.35,
            sc.fillna(0.0) < 0.35,
            tm.fillna(0.0) < 0.10,
            meta,
        )

    if n < 25:
        meta.update(
            {
                "note": f"Only {n} rows — using pipeline-style floors.",
                "topic_thr": LOW_TOPIC_CONFIDENCE,
                "sent_thr": LOW_SENTIMENT_CONFIDENCE,
                "margin_thr": 0.10,
            }
        )
        return (
            tc.fillna(0.0) < LOW_TOPIC_CONFIDENCE,
            sc.fillna(0.0) < LOW_SENTIMENT_CONFIDENCE,
            tm.fillna(0.0) < 0.10,
            meta,
        )

    if mode == "percentile":
        q = max(0.03, min(0.35, tail_pct / 100.0))
        t_thr = float(tc.quantile(q))
        s_thr = float(sc.quantile(q))
        m_thr = float(tm.quantile(q))
        meta.update({"topic_thr": t_thr, "sent_thr": s_thr, "margin_thr": m_thr, "tail_pct": tail_pct, "quantile": q})
        return (
            tc.fillna(0.0) < t_thr,
            sc.fillna(0.0) < s_thr,
            tm.fillna(0.0) < m_thr,
            meta,
        )

    # z_score: flag if metric is unusually low vs corpus (mean − z·std), never stricter than fixed 0.35.
    z = max(0.5, min(3.0, float(z_mult)))
    meta["z"] = z

    def z_threshold(series: pd.Series, hard_cap: float) -> float:
        s = series.dropna()
        if len(s) < 25:
            return float(hard_cap)
        mu = float(s.mean())
        sig = float(s.std(ddof=0))
        if sig < 0.015:
            thr = float(s.quantile(max(0.05, tail_pct / 100.0)))
        else:
            thr = mu - z * sig
        thr = min(float(hard_cap), thr)
        thr = max(float(s.quantile(0.02)), thr)
        return thr

    t_thr = z_threshold(tc, 0.40)
    s_thr = z_threshold(sc, 0.40)
    m_thr = z_threshold(tm, 0.12)
    meta.update({"topic_thr": t_thr, "sent_thr": s_thr, "margin_thr": m_thr})
    return (
        tc.fillna(0.0) < t_thr,
        sc.fillna(0.0) < s_thr,
        tm.fillna(0.0) < m_thr,
        meta,
    )


def enrich_topic_borderline_and_scores(df: pd.DataFrame) -> pd.DataFrame:
    """
    Backfill topic_confidence / topic_margin when missing, derive runner-up topic and
    topic_near_tie for UI (matches sentimentanalysis when those fields are absent).
    """
    out = df.copy()
    topic_order = list(TOPIC_DEFINITIONS.keys())
    def _num_series(col: str, default: float = np.nan) -> pd.Series:
        if col in out.columns:
            return pd.to_numeric(out[col], errors="coerce")
        return pd.Series(default, index=out.index, dtype="float64")

    ts = _num_series("topic_score")
    if "topic_confidence" not in out.columns:
        out["topic_confidence"] = ts
    else:
        out["topic_confidence"] = pd.to_numeric(out["topic_confidence"], errors="coerce").fillna(ts)

    def _margin_and_second(hy):
        if not isinstance(hy, list) or len(hy) < 2:
            return pd.Series({"topic_margin": np.nan, "topic_second_best": None})
        arr = np.asarray(hy, dtype=float)
        if len(arr) != len(topic_order):
            idx = np.argsort(arr)
            j1, j2 = int(idx[-1]), int(idx[-2])
            return pd.Series(
                {
                    "topic_margin": float(arr[j1] - arr[j2]),
                    "topic_second_best": None,
                }
            )
        idx = np.argsort(arr)
        j1, j2 = int(idx[-1]), int(idx[-2])
        return pd.Series(
            {
                "topic_margin": float(arr[j1] - arr[j2]),
                "topic_second_best": topic_order[j2],
            }
        )

    if "topic_scores_hybrid" in out.columns:
        ms = out["topic_scores_hybrid"].apply(_margin_and_second)
        out["topic_margin"] = _num_series("topic_margin")
        out["topic_margin"] = out["topic_margin"].fillna(ms["topic_margin"])
        if "topic_second_best" not in out.columns:
            out["topic_second_best"] = ms["topic_second_best"]
        else:
            mask = out["topic_second_best"].isna() | (out["topic_second_best"].astype(str) == "")
            out.loc[mask, "topic_second_best"] = ms.loc[mask, "topic_second_best"]
    else:
        if "topic_margin" not in out.columns:
            out["topic_margin"] = np.nan
        else:
            out["topic_margin"] = pd.to_numeric(out["topic_margin"], errors="coerce")
        if "topic_second_best" not in out.columns:
            out["topic_second_best"] = None

    med_m = float(pd.to_numeric(out["topic_margin"], errors="coerce").median())
    fb = med_m if not np.isnan(med_m) else TOPIC_MARGIN_DRIFT
    out["topic_margin"] = pd.to_numeric(out["topic_margin"], errors="coerce").fillna(fb)

    tm = pd.to_numeric(out["topic_margin"], errors="coerce")
    if "topic_near_tie" not in out.columns:
        out["topic_near_tie"] = tm < TOPIC_NEAR_TIE_MARGIN
    else:
        ntt = out["topic_near_tie"].fillna(False).astype(bool)
        out["topic_near_tie"] = ntt | (tm < TOPIC_NEAR_TIE_MARGIN)

    score_series = _num_series("score")
    if "sentiment_confidence" not in out.columns:
        out["sentiment_confidence"] = score_series.fillna(0.0)
    else:
        out["sentiment_confidence"] = pd.to_numeric(out["sentiment_confidence"], errors="coerce").fillna(
            score_series
        ).fillna(0.0)

    return out


def render_outlier_review_page(df_sent: pd.DataFrame):
    st.header("Outlier Review Queue")
    st.write("Review low-confidence, drift-risk, and manually overridden sentences. Apply overrides directly from this queue.")

    st.markdown(
        "Topic and sentiment confidences are **probabilities between 0 and 1** (often **below 0.35** even when the label is fine) "
        "because the pipeline blends **embedding similarity + NLI** over eight topics, or **three sentiment models**, "
        "then normalises. That is expected; the old queue used **fixed 0.35 / 0.10 cutoffs**, which treated typical scores as outliers. "
        "Use **corpus-relative** thresholds below unless you intentionally want a strict absolute screen."
    )
    st.caption(
        f"**Topic borderline:** margin (top hybrid minus runner-up) below **{TOPIC_NEAR_TIE_MARGIN:.2f}** — "
        f"highlighted in the table; optionally add all such rows to the queue below. "
        f"Pipeline drift flag uses margin **{TOPIC_MARGIN_DRIFT:.2f}**."
    )

    df = enrich_topic_borderline_and_scores(df_sent.copy())
    if "needs_review" not in df.columns:
        df["needs_review"] = False
    if "review_reasons" not in df.columns:
        df["review_reasons"] = [[] for _ in range(len(df))]
    if "manual_topic_override_applied" not in df.columns:
        df["manual_topic_override_applied"] = False
    if "manual_sentiment_override_applied" not in df.columns:
        df["manual_sentiment_override_applied"] = False

    with st.expander("Queue threshold mode", expanded=True):
        mode = st.radio(
            "Soft rules (topic / sentiment / margin)",
            options=["percentile", "z_score", "fixed"],
            index=0,
            horizontal=True,
            help="Soft rules are combined with pipeline needs_review and overrides. "
            "Percentile: flag worst tail of the corpus per metric. "
            "Z-score: flag values below mean − z·std (with safeguards). "
            "Fixed: legacy 0.35 / 0.35 / 0.10 (often flags almost all rows).",
        )
        tail_pct = 12
        z_mult = 1.5
        if mode == "percentile":
            tail_pct = st.slider("Approx. bottom percentile per metric (%)", 5, 30, 12)
        elif mode == "z_score":
            z_mult = st.slider("Z below mean (topic, sentiment, margin)", 0.8, 2.5, 1.5, step=0.1)
            tail_pct = st.slider("Fallback percentile if spread is flat (%)", 5, 20, 10)
        hide_leadership_only = st.checkbox(
            "Hide rows that only mention leadership keywords (no other review reason)",
            value=False,
            help="Sentences matching CEO/chair/leadership etc. are flagged for governance review; "
            "exclude them here if they crowd out real low-confidence cases.",
        )
        include_topic_near_ties = st.checkbox(
            "Also include topic borderline rows (near tie between top two topics)",
            value=False,
            help=f"Adds sentences where the hybrid margin is below {TOPIC_NEAR_TIE_MARGIN} "
            "(assigned topic and runner-up are similarly likely), even if no other rule matched.",
        )

    low_topic, low_sent, low_margin, thr_meta = compute_soft_outlier_masks(
        df, mode=mode, tail_pct=tail_pct, z_mult=z_mult
    )

    review_flag = df["review_reasons"].apply(lambda x: isinstance(x, list) and len(x) > 0)
    pipeline_flag = df["needs_review"] == True
    soft_union = low_topic | low_sent | low_margin
    override_union = (df["manual_topic_override_applied"] == True) | (df["manual_sentiment_override_applied"] == True)

    leadership_only = df["review_reasons"].apply(
        lambda r: _leadership_keyword_only(_review_reasons_as_list(r))
    )

    near_tie_union = df["topic_near_tie"].fillna(False).astype(bool)
    include_row = pipeline_flag | review_flag | soft_union | override_union
    if include_topic_near_ties:
        include_row = include_row | near_tie_union
    if hide_leadership_only:
        include_row = include_row & ~(leadership_only & ~soft_union & ~override_union)

    flagged = df.loc[include_row].copy()

    if flagged.empty:
        st.success("No flagged sentences currently in review queue.")
    else:
        flagged["review_reasons"] = flagged["review_reasons"].apply(
            lambda x: ", ".join(x) if isinstance(x, list) else str(x)
        )
        flagged = flagged.sort_values(
            by=["topic_confidence", "sentiment_confidence", "topic_margin"],
            ascending=[True, True, True],
        )
        st.write(f"Flagged rows: {len(flagged)}")
        thr_line = (
            f"Soft thresholds ({thr_meta.get('mode', mode)}): "
            f"topic < {float(thr_meta['topic_thr']):.3f} | "
            f"sentiment < {float(thr_meta['sent_thr']):.3f} | "
            f"margin < {float(thr_meta['margin_thr']):.3f}"
        )
        if "note" in thr_meta:
            thr_line += f" — {thr_meta['note']}"
        st.caption(thr_line)
        st.caption(
            f"Counts — pipeline needs_review: {int(pipeline_flag.sum())} | "
            f"non-empty review_reasons: {int(review_flag.sum())} | "
            f"soft low topic: {int(low_topic.sum())} | soft low sentiment: {int(low_sent.sum())} | "
            f"soft low margin: {int(low_margin.sum())} | "
            f"leadership-keyword-only: {int(leadership_only.sum())} | "
            f"topic borderline (margin < {TOPIC_NEAR_TIE_MARGIN}): {int(near_tie_union.sum())} in corpus, "
            f"{int(flagged['topic_near_tie'].fillna(False).sum())} in this table"
        )
        display_cols = [
            "global_index",
            "article_filename",
            "sentence",
            "topic_name",
            "topic_second_best",
            "topic_bucket",
            "sentiment_display",
            "topic_score",
            "topic_confidence",
            "sentiment_confidence",
            "topic_margin",
            "topic_near_tie",
            "review_reasons",
            "manual_sentiment_override_applied",
            "manual_topic_override_applied",
        ]
        display_cols = [c for c in display_cols if c in flagged.columns]
        view = flagged[display_cols].head(250).copy()
        view["topic_near_tie"] = view["topic_near_tie"].fillna(False).astype(bool)

        def _highlight_near_tie(row: pd.Series) -> list:
            if bool(row.get("topic_near_tie")):
                return ["background-color: #fff8e1; border-left: 3px solid #b8860b"] * len(row)
            return [""] * len(row)

        try:
            styled = (
                view.style
                .apply(_highlight_near_tie, axis=1)
                .set_properties(**{"color": SECONDARY_BLUE, "font-family": "Garamond"})
            )
            st.dataframe(styled, use_container_width=True)
        except Exception:
            st.dataframe(view, use_container_width=True)

    st.subheader("Apply Manual Override From Queue")
    default_idx = int(flagged["global_index"].iloc[0]) if not flagged.empty and "global_index" in flagged.columns else 0
    render_manual_override_controls(df_sent, key_prefix="review_queue", default_global_index=default_idx)

    st.markdown("### Override Maintenance")
    wipe_confirm = st.checkbox(
        "Confirm: wipe all manual sentiment/topic overrides before next pipeline refresh",
        value=False,
        key="wipe_overrides_confirm",
    )
    if st.button("Wipe all overrides", key="wipe_overrides_btn", type="secondary"):
        if not wipe_confirm:
            st.warning("Tick the confirmation checkbox first to wipe overrides.")
        else:
            ok, message = clear_all_overrides_in_master()
            if ok:
                st.success(f"✓ {message}")
                st.info("Overrides are now cleared in `master.json`. The next pipeline run will start from a clean override state.")
                st.rerun()
            else:
                st.error(f"✗ {message}")


def apply_override_to_master(global_index: int, sentiment_override: str, topic_override: str):
    try:
        master = json.loads(MASTER_JSON.read_text(encoding="utf-8"))
        sent_override = None if sentiment_override == "(no override)" else sentiment_override
        topic_override_val = None if topic_override == "(no override)" else topic_override

        if update_sentence_override(master, global_index, sent_override, topic_override_val):
            MASTER_JSON.write_text(json.dumps(master, indent=2), encoding="utf-8")
            load_master.clear()
            return True, f"Override applied to global_index {global_index}."
        return False, f"Global index {global_index} not found in master.json."
    except Exception as exc:
        return False, f"Failed to apply override: {exc}"


def clear_all_overrides_in_master():
    try:
        master = json.loads(MASTER_JSON.read_text(encoding="utf-8"))
        sentences = master.get("sentences", [])
        if not isinstance(sentences, list):
            return False, "master.json has no valid `sentences` list."

        changed = 0
        for s in sentences:
            if not isinstance(s, dict):
                continue
            had_override = bool(
                s.get("manual_sentiment_override") is not None
                or s.get("manual_topic_override") is not None
                or s.get("manual_sentiment_override_applied") is True
                or s.get("manual_topic_override_applied") is True
            )
            s["manual_sentiment_override"] = None
            s["manual_topic_override"] = None
            s["manual_sentiment_override_applied"] = False
            s["manual_topic_override_applied"] = False
            if had_override:
                changed += 1

        MASTER_JSON.write_text(json.dumps(master, indent=2), encoding="utf-8")
        load_master.clear()
        return True, f"Cleared overrides on {changed} sentence row(s)."
    except Exception as exc:
        return False, f"Failed to clear overrides: {exc}"


def render_manual_override_controls(df_sent: pd.DataFrame, key_prefix: str = "override", default_global_index: int = 0):
    col1, col2, col3 = st.columns(3)
    with col1:
        global_index = st.number_input(
            "Global Index to Override",
            min_value=0,
            step=1,
            value=int(default_global_index),
            key=f"{key_prefix}_global_index",
        )
    with col2:
        sentiment_override = st.selectbox(
            "Override Sentiment",
            ["(no override)", "very_negative", "negative", "neutral", "positive", "very_positive"],
            key=f"{key_prefix}_sentiment",
        )
    with col3:
        topic_override = st.selectbox(
            "Override Topic",
            ["(no override)", "None"] + list(TOPIC_DEFINITIONS.keys()),
            key=f"{key_prefix}_topic",
        )

    if st.button("Apply Override to Master", key=f"{key_prefix}_apply_btn"):
        ok, message = apply_override_to_master(int(global_index), sentiment_override, topic_override)
        if ok:
            st.success(f"✓ {message}")
            st.info(
                "Dashboard charts and storyboard exports read from `master.json` after reload—no local pipeline rerun needed for overrides. "
                "For new/changed articles or full model refreshes, run the **Run Sentiment Pipeline** GitHub Action."
            )
            st.rerun()
        else:
            st.error(f"✗ {message}")

# ---------------------------------------------------------
# UI WRAPPERS FOR PPTX & NARRATIVE (Point 4: Implementation Support)
# ---------------------------------------------------------

def render_powerpoint_storyboard(df_sent, df_topics, bucket_sizes):
    st.header("PowerPoint Storyboard Export")
    st.write("""
    Generate a professional multi-slide PowerPoint storyboard from the current dataset.
    Charts are cached in `./powerpoint/` and reused across exports.
    Design: neutral white theme with Helvetica/Garamond typography.
    """)

    with st.expander("Export diagnostics", expanded=False):
        st.caption("Use this to compare local vs hosted rendering environments.")
        try:
            altair_version = getattr(alt, "__version__", "unknown")
        except Exception:
            altair_version = "unknown"
        try:
            vl_convert_version = getattr(vlc, "__version__", "installed (version unknown)") if vlc is not None else "not installed"
        except Exception:
            vl_convert_version = "installed (version lookup failed)"
        st.code(
            "\n".join(
                [
                    f"python_version={sys.version.split()[0]}",
                    f"streamlit_version={st.__version__}",
                    f"altair_version={altair_version}",
                    f"vl_convert_available={vlc is not None}",
                    f"vl_convert_version={vl_convert_version}",
                    f"export_scale_factor={EXPORT_SCALE_FACTOR}",
                    f"chart_export_dir={CHART_EXPORT_DIR.resolve()}",
                ]
            ),
            language="text",
        )
    
    if st.button("Generate Storyboard PPTX", key="gen_pptx_btn"):
        with st.spinner("Building storyboard (this may take a minute)..."):
            try:
                df_article_sent = compute_article_sentiment(df_sent)
                df_polarity = compute_bucket_polarity(df_sent)
                slides = build_storyboard_slides(df_sent, df_article_sent, df_polarity, bucket_sizes)
                pptx_bytes = export_storyboard_to_pptx(slides)
                st.success(f"✓ Storyboard built — {len(slides)} slides")
                render_export_png_preview(slides)
                st.download_button(
                    "📥 Download PowerPoint",
                    data=pptx_bytes,
                    file_name=f"storyboard_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                st.error(f"Failed to build storyboard: {e}")


def render_export_png_preview(slides: list):
    image_rows = []
    for slide_idx, slide in enumerate(slides or [], start=1):
        for image_path in slide.get("images", []) or []:
            p = Path(str(image_path))
            if p.exists():
                image_rows.append((slide_idx, slide.get("title", f"Slide {slide_idx}"), p))

    if not image_rows:
        return

    with st.expander("Export PNG preview", expanded=True):
        st.caption("These are the exact PNG files inserted into the PowerPoint download.")
        for i in range(0, len(image_rows), 2):
            cols = st.columns(2)
            for col, row in zip(cols, image_rows[i:i + 2]):
                slide_idx, slide_title, image_path = row
                with col.container():
                    st.caption(f"Slide {slide_idx}: {slide_title}")
                    st.image(str(image_path), use_container_width=True)

def render_narrative_export(df_sent):
    st.header("Narrative Export")
    st.write("Generate a plain-text narrative summary and download as .txt for distribution.")
    
    if st.button("Generate Narrative", key="gen_narrative_btn"):
        try:
            text = build_narrative(df_sent)
            st.code(text, language="plaintext")
            st.download_button(
                "📥 Download Narrative (.txt)",
                data=text,
                file_name=f"narrative_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
            )
        except Exception as e:
            st.error(f"Failed to generate narrative: {e}")

# ---------------------------------------------------------
# ... REST OF EXISTING FUNCTIONS REMAIN UNCHANGED ...
# ---------------------------------------------------------

def apply_topic_buckets(df_sent, df_topics):
    df_sent = df_sent.copy()
    df_topics = df_topics.copy()

    if "topic_name" not in df_sent.columns:
        df_sent["topic_name"] = "None"
    if "topic_name" not in df_topics.columns:
        df_topics["topic_name"] = "None"

    df_sent["topic_bucket"] = df_sent["topic_name"].map(TOPIC_BUCKET_MAP).fillna("None")
    df_topics["topic_bucket"] = df_topics["topic_name"].map(TOPIC_BUCKET_MAP).fillna("None")

    if "sentence" not in df_sent.columns:
        df_sent["sentence"] = df_sent.get("text", "")

    bucket_sizes = (
        df_sent[df_sent["topic_bucket"].ne("None")]
        .groupby("topic_bucket")["sentence"]
        .count()
        .reset_index()
        .rename(columns={"sentence": "size"})
        .sort_values("size", ascending=False)
    )

    return df_sent, df_topics, bucket_sizes

# ... (ALL REMAINING FUNCTIONS FROM ORIGINAL FILE COPIED HERE UNCHANGED) ...

def compute_bucket_polarity(df_sent: pd.DataFrame) -> pd.DataFrame:
    """Strategic-bucket polarity: weighted mix of sentence display sentiments in each topic_bucket (not topic-model polarity)."""
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return pd.DataFrame()

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    rows = []
    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]

        polarity = sum(
            SENTIMENT_WEIGHTS[row["sentiment_display"]] * row["percent"]
            for _, row in sub.iterrows()
        )

        pos = sub[sub["sentiment_display"].isin(["Positive", "Very Positive"])]["percent"].sum()
        neg = sub[sub["sentiment_display"].isin(["Negative", "Very Negative"])]["percent"].sum()
        neu = sub[sub["sentiment_display"].eq("Neutral")]["percent"].sum()

        rows.append({
            "topic_bucket": bucket,
            "bucket_short": BUCKET_SHORT.get(bucket, bucket),
            "polarity": polarity,
            "positive_percent": pos,
            "negative_percent": neg,
            "neutral_percent": neu,
        })

    return pd.DataFrame(rows)


def compute_article_bucket_polarity(df_sent: pd.DataFrame) -> pd.DataFrame:
    """Per-bucket polarity when each article is summarized to one tone (article-level roll-up)."""
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty or "article_id" not in df.columns:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS).fillna(0.0)
    df["weighted"] = df["sentiment_weight"] * df.get("topic_score", 1.0)

    art_scores = (
        df.groupby("article_id")["weighted"].sum().reset_index().rename(columns={"weighted": "article_score"})
    )

    rep_bucket = (
        df.groupby(["article_id", "topic_bucket"]).size().reset_index(name="count")
        .sort_values(["article_id", "count"], ascending=[True, False])
        .groupby("article_id").first().reset_index()[["article_id", "topic_bucket"]]
    )

    art = art_scores.merge(rep_bucket, on="article_id", how="left")

    def label_from_score(s):
        if s > 0:
            return "Positive"
        if s < 0:
            return "Negative"
        return "Neutral"

    art["article_sentiment"] = art["article_score"].apply(label_from_score)

    pivot = (
        art.groupby(["topic_bucket", "article_sentiment"]).size().reset_index(name="count")
    )
    totals = (
        pivot.groupby("topic_bucket")["count"].sum().reset_index().rename(columns={"count": "total_articles"})
    )
    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total_articles"]) * 100

    rows = []
    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]
        pos = sub[sub["article_sentiment"] == "Positive"]["percent"].sum()
        neg = sub[sub["article_sentiment"] == "Negative"]["percent"].sum()
        neu = sub[sub["article_sentiment"] == "Neutral"]["percent"].sum()
        total_articles = int(sub["total_articles"].iloc[0]) if not sub.empty else 0

        rows.append({
            "topic_bucket": bucket,
            "bucket_short": BUCKET_SHORT.get(bucket, bucket),
            "article_positive_percent": pos,
            "article_negative_percent": neg,
            "article_neutral_percent": neu,
            "total_articles": total_articles,
        })

    return pd.DataFrame(rows)


def generate_bucket_summary(df_sent: pd.DataFrame) -> dict:
    summaries = {}

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return summaries

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100

    for bucket in pivot["topic_bucket"].unique():
        sub = pivot[pivot["topic_bucket"] == bucket]

        polarity = sum(
            SENTIMENT_WEIGHTS[row["sentiment_display"]] * row["percent"]
            for _, row in sub.iterrows()
        )

        pos = sub[sub["sentiment_display"].isin(["Positive", "Very Positive"])]["percent"].sum()
        neg = sub[sub["sentiment_display"].isin(["Negative", "Very Negative"])]["percent"].sum()
        neu = sub[sub["sentiment_display"].eq("Neutral")]["percent"].sum()

        if polarity > 0:
            tone = "strong positive sentiment"
            arrow = "↑"
        elif polarity < 0:
            tone = "strong negative sentiment"
            arrow = "↓"
        else:
            tone = "mixed or neutral sentiment"
            arrow = "→"

        summaries[bucket] = (
            f"{arrow} {bucket} shows {tone}. "
            f"Positive coverage: {pos:.1f}%. "
            f"Negative coverage: {neg:.1f}%. "
            f"Neutral coverage: {neu:.1f}%. "
            f"Polarity score: {polarity:.1f}."
        )

    return summaries


def get_sentiment_drivers(df_sent: pd.DataFrame, top_n: int = 5) -> dict:
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return {}

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
    df["driver_score"] = df["sentiment_weight"] * df["topic_score"]

    drivers = {}

    for bucket in df["topic_bucket"].unique():
        sub = df[df["topic_bucket"] == bucket]

        pos = (
            sub[sub["sentiment_weight"] > 0]
            .sort_values("driver_score", ascending=False)
            .head(top_n)
        )
        neg = (
            sub[sub["sentiment_weight"] < 0]
            .sort_values("driver_score")
            .head(top_n)
        )

        drivers[bucket] = {
            "positive": pos[["sentence", "sentiment_display", "topic_score", "driver_score"]],
            "negative": neg[["sentence", "sentiment_display", "topic_score", "driver_score"]],
        }

    return drivers


def get_global_sentiment_drivers(df_sent: pd.DataFrame, top_n: int = 3):
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()

    if df.empty:
        return pd.DataFrame(), pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)
    df["driver_score"] = df["sentiment_weight"] * df["topic_score"]

    pos = (
        df[df["sentiment_weight"] > 0]
        .sort_values("driver_score", ascending=False)
        .head(top_n)
    )
    neg = (
        df[df["sentiment_weight"] < 0]
        .sort_values("driver_score")
        .head(top_n)
    )

    return pos, neg


def compute_article_sentiment(df_sent: pd.DataFrame) -> pd.DataFrame:
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "article_id" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent.copy()
    df = df[df["article_id"].notna()]
    if df.empty:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)

    agg = (
        df.groupby("article_id")
        .agg(
            avg_weight=("sentiment_weight", "mean"),
            sentence_count=("sentence", "size"),
        )
        .reset_index()
    )

    def classify_tone(weight: float) -> str:
        if weight > 0.1:
            return "Positive"
        if weight < -0.1:
            return "Negative"
        return "Neutral"

    agg["overall_tone"] = agg["avg_weight"].apply(classify_tone)
    return agg


def build_sentence_distribution_chart(df_sent: pd.DataFrame):
    if df_sent is None or df_sent.empty:
        return None

    df = (
        df_sent.groupby("sentiment_display")
        .size()
        .reset_index(name="count")
    )

    df["sentiment_display"] = pd.Categorical(
        df["sentiment_display"],
        categories=SENTIMENT_ORDER,
        ordered=True,
    )
    df = df.sort_values("sentiment_display")

    color_scale = alt.Scale(
        domain=SENTIMENT_ORDER,
        range=[SENTIMENT_COLORS[s] for s in SENTIMENT_ORDER],
    )

    bars = (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(title="Sentiment", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            y=alt.Y(
                "count:Q",
                axis=alt.Axis(title="Number of sentences", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            color=alt.Color(
                "sentiment_display:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=["sentiment_display", "count"],
        )
    )

    chart = bars.properties(
        title="Sentence Sentiment Distribution",
        width=DEFAULT_CHART_WIDTH,
        height=DEFAULT_CHART_HEIGHT,
    )

    return chart


def build_article_tone_chart(df_article_sent: pd.DataFrame):
    if df_article_sent is None or df_article_sent.empty:
        return None

    df = (
        df_article_sent.groupby("overall_tone")
        .size()
        .reset_index(name="count")
    )

    tone_order = ["Negative", "Neutral", "Positive"]
    df["overall_tone"] = pd.Categorical(
        df["overall_tone"],
        categories=tone_order,
        ordered=True,
    )
    df = df.sort_values("overall_tone")

    tone_colors = {
        "Negative": SENTIMENT_COLORS["Negative"],
        "Neutral": SENTIMENT_COLORS["Neutral"],
        "Positive": SENTIMENT_COLORS["Positive"],
    }

    color_scale = alt.Scale(
        domain=tone_order,
        range=[tone_colors[t] for t in tone_order],
    )

    bars = (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X(
                "overall_tone:N",
                sort=tone_order,
                axis=alt.Axis(title="Article tone", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            y=alt.Y(
                "count:Q",
                axis=alt.Axis(title="Number of articles", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            color=alt.Color(
                "overall_tone:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=["overall_tone", "count"],
        )
    )

    chart = bars.properties(
        title="Article Tone Distribution",
        width=DEFAULT_CHART_WIDTH,
        height=DEFAULT_CHART_HEIGHT,
    )

    return chart


def build_bucket_polarity_bar_chart(df_polarity: pd.DataFrame):
    if df_polarity is None or df_polarity.empty:
        return None

    df = df_polarity.copy()
    if "bucket_short" not in df.columns:
        df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)

    bucket_domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]

    chart = (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X(
                "bucket_short:N",
                sort=bucket_domain,
                axis=alt.Axis(title="Bucket (code)", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            y=alt.Y(
                "polarity:Q",
                axis=alt.Axis(
                    title="Polarity score",
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                    offset=8,
                    tickSize=6,
                ),
            ),
            color=alt.condition(
                alt.datum.polarity > 0,
                alt.value("#2ca02c"),
                alt.value("#d62728"),
            ),
            tooltip=[
                "topic_bucket",
                "bucket_short",
                "polarity",
                "positive_percent",
                "negative_percent",
                "neutral_percent",
            ],
        )
    ).properties(
        title="Bucket Polarity",
        width=DEFAULT_CHART_WIDTH,
        height=DEFAULT_CHART_HEIGHT,
    )

    return chart


def build_bucket_sizes_chart(bucket_sizes: pd.DataFrame):
    if bucket_sizes is None or bucket_sizes.empty:
        return None
    if "topic_bucket" not in bucket_sizes.columns or "size" not in bucket_sizes.columns:
        return None

    df = bucket_sizes[bucket_sizes["topic_bucket"].ne("None")].copy()
    if df.empty:
        return None

    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    bucket_domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]

    bars = (
        alt.Chart(df)
        .mark_bar(color=SECONDARY_BLUE)
        .encode(
            x=alt.X(
                "bucket_short:N",
                sort=bucket_domain,
                axis=alt.Axis(title="Bucket (code)", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            y=alt.Y(
                "size:Q",
                axis=alt.Axis(title="Sentences", labelPadding=AXIS_LABEL_PADDING, titlePadding=AXIS_TITLE_PADDING),
            ),
            tooltip=["topic_bucket", "bucket_short", "size"],
        )
    )

    chart = bars.properties(
        title="Bucket Sizes",
        width=DEFAULT_CHART_WIDTH,
        height=DEFAULT_CHART_HEIGHT,
    )

    return chart


def build_bucket_balance_table(df_sent: pd.DataFrame) -> pd.DataFrame:
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "topic_bucket" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]
    if df.empty:
        return pd.DataFrame()

    df["sentiment_weight"] = df["sentiment_display"].map(SENTIMENT_WEIGHTS)

    is_pos = df["sentiment_display"].isin(["Positive", "Very Positive"])
    is_neg = df["sentiment_display"].isin(["Negative", "Very Negative"])

    agg = (
        df.groupby("topic_bucket")
        .agg(
            positive_count=("sentiment_display", lambda s: int(is_pos.loc[s.index].sum())),
            negative_count=("sentiment_display", lambda s: int(is_neg.loc[s.index].sum())),
            total_count=("sentiment_display", "size"),
            avg_intensity=("sentiment_weight", "mean"),
        )
        .reset_index()
    )

    agg["net_balance"] = agg["positive_count"] - agg["negative_count"]
    agg["bucket_short"] = agg["topic_bucket"].map(BUCKET_SHORT)

    return agg


def bucket_balance_bubble(df_sent: pd.DataFrame):
    table = build_bucket_balance_table(df_sent)
    if table is None or table.empty:
        return None

    table = table.sort_values("net_balance")
    if len(table) > 1:
        offsets = np.linspace(-0.03, 0.03, len(table))
        table["label_y"] = table["avg_intensity"] + offsets
    else:
        table["label_y"] = table["avg_intensity"]
    base = alt.Chart(table).encode(
        x=alt.X(
            "net_balance:Q",
            axis=alt.Axis(
                title="Net balance",
                offset=8,
                labelPadding=AXIS_LABEL_PADDING,
                titlePadding=AXIS_TITLE_PADDING,
            ),
            scale=alt.Scale(nice=True, padding=36),
        ),
        y=alt.Y(
            "avg_intensity:Q",
            axis=alt.Axis(
                title="Intensity",
                offset=8,
                labelPadding=AXIS_LABEL_PADDING,
                titlePadding=AXIS_TITLE_PADDING,
            ),
            scale=alt.Scale(nice=True, padding=18),
        ),
    )

    bubbles = base.mark_circle(opacity=0.7, stroke="black", strokeWidth=0.5).encode(
        size=alt.Size(
            "total_count:Q",
            title="Total sentences in bucket",
            scale=alt.Scale(range=[200, 2000]),
        ),
        color=alt.Color(
            "avg_intensity:Q",
            title="Average intensity",
            scale=alt.Scale(scheme="redyellowgreen"),
            legend=None,
        ),
        tooltip=[
            "topic_bucket",
            "net_balance",
            "avg_intensity",
            "positive_count",
            "negative_count",
            "total_count",
        ],
    )

    labels = (
        base.mark_text(
            baseline="middle",
            align="left",
            dx=10,
            fontSize=12,
            color=PRIMARY_BLUE,
            strokeWidth=0.8,
        ).encode(
            text=alt.Text("bucket_short:N", title=None),
            y=alt.Y("label_y:Q"),
        )
    )

    vline = alt.Chart(table).mark_rule(color="#555555", strokeWidth=2).encode(x=alt.datum(0))
    hline = alt.Chart(table).mark_rule(color="#555555", strokeWidth=2).encode(y=alt.datum(0))

    return (vline + hline + bubbles + labels).properties(
        title="Bucket Balance Map",
        width=DEFAULT_CHART_WIDTH,
        height=DEFAULT_CHART_HEIGHT,
    )


def bucket_sentiment_heatmap(df_sent: pd.DataFrame):
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return None

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100
    pivot["bucket_short"] = pivot["topic_bucket"].map(BUCKET_SHORT)

    chart = (
        alt.Chart(pivot)
        .mark_rect()
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(
                    title="Sentiment",
                    labelLimit=AXIS_TITLE_LIMIT,
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            y=alt.Y(
                "bucket_short:N",
                sort=[BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT],
                axis=alt.Axis(
                    title="Bucket (code)",
                    labelLimit=1000,
                    labelAlign="right",
                    labelBaseline="middle",
                    labelAngle=0,
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            color=alt.Color(
                "percent:Q",
                scale=alt.Scale(scheme="blues"),
                title="Percent of Bucket",
            ),
            tooltip=[
                "topic_bucket",
                "sentiment_display",
                "percent",
                "count",
                "total",
            ],
        ).properties(
            title="Bucket x Sentiment Composition",
            width=DEFAULT_CHART_WIDTH,
            height=DEFAULT_CHART_HEIGHT,
        )
    )

    return chart


def bucket_sentiment_bubble(df_sent: pd.DataFrame):
    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return None

    df["count"] = 1

    pivot = (
        df.groupby(["topic_bucket", "sentiment_display"])["count"]
        .sum()
        .reset_index()
    )

    totals = (
        pivot.groupby("topic_bucket")["count"]
        .sum()
        .reset_index()
        .rename(columns={"count": "total"})
    )

    pivot = pivot.merge(totals, on="topic_bucket")
    pivot["percent"] = (pivot["count"] / pivot["total"]) * 100
    pivot["bucket_short"] = pivot["topic_bucket"].map(BUCKET_SHORT)

    color_scale = alt.Scale(
        domain=SENTIMENT_ORDER,
        range=[SENTIMENT_COLORS[s] for s in SENTIMENT_ORDER],
    )

    chart = (
        alt.Chart(pivot)
        .mark_circle(opacity=0.85, stroke="black", strokeWidth=0.2)
        .encode(
            x=alt.X(
                "sentiment_display:N",
                sort=SENTIMENT_ORDER,
                axis=alt.Axis(
                    title="Sentiment",
                    labelLimit=AXIS_TITLE_LIMIT,
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            y=alt.Y(
                "bucket_short:N",
                sort=[BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT],
                axis=alt.Axis(
                    title="Bucket (code)",
                    labelLimit=AXIS_TITLE_LIMIT,
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            size=alt.Size(
                "percent:Q",
                title="Percent of Bucket",
                scale=alt.Scale(range=[50, 1000]),
            ),
            color=alt.Color(
                "sentiment_display:N",
                scale=color_scale,
                legend=None,
            ),
            tooltip=[
                "topic_bucket",
                "sentiment_display",
                "percent",
                "count",
                "total",
            ],
        ).properties(
            title="Bucket x Sentiment Bubble View",
            width=DEFAULT_CHART_WIDTH,
            height=DEFAULT_CHART_HEIGHT,
        )
    )

    return chart


def topic_drift_heatmap(df_sent: pd.DataFrame):
    """
    For each governance bucket (column), show each fine-topic code (row) as % of all
    sentences in that bucket. Multiple pipeline topic_name strings that share the same
    axis code (e.g. Leadership vs Leadership & Governance → LG) are summed into one cell.
    """
    if "topic_bucket" not in df_sent.columns or "topic_name" not in df_sent.columns:
        return None

    df = df_sent.copy()
    df = df[df["topic_bucket"].ne("None")]
    df = df[df["topic_bucket"].ne("Other")]
    df = df[df["topic_name"].notna()]

    if df.empty:
        return None

    df["bucket_short"] = df["topic_bucket"].map(BUCKET_SHORT)
    df = df[df["bucket_short"].notna()]
    if df.empty:
        return None

    df["topic_graph"] = (
        df["topic_name"].map(TOPIC_GRAPH_LABEL).fillna(df["topic_name"].astype(str).str.slice(0, 24))
    )
    df["count"] = 1

    bucket_totals = (
        df.groupby("bucket_short", as_index=False)["count"]
        .sum()
        .rename(columns={"count": "bucket_total"})
    )

    pivot = (
        df.groupby(["topic_graph", "bucket_short"], as_index=False)
        .agg(
            count=("count", "sum"),
            topic_names=(
                "topic_name",
                lambda s: ", ".join(sorted(pd.Series(s).dropna().unique().astype(str))),
            ),
        )
    )
    pivot = pivot.merge(bucket_totals, on="bucket_short", how="left")
    bt = pivot["bucket_total"].replace(0, np.nan)
    pivot["percent"] = ((pivot["count"] / bt) * 100).fillna(0.0)
    pivot["topic_bucket"] = pivot["bucket_short"].map(BUCKET_SHORT_TO_NAME)

    bucket_domain = [BUCKET_SHORT[b] for b in BUCKET_ORDER if b in BUCKET_SHORT]
    known_labels = [TOPIC_GRAPH_LABEL[t] for t in TOPIC_DEFINITIONS.keys() if t in TOPIC_GRAPH_LABEL]
    present = set(pivot["topic_graph"].astype(str))
    topic_domain = [lab for lab in known_labels if lab in present] + sorted(
        present.difference(known_labels), key=str
    )
    if not topic_domain:
        return None

    pct_max = float(pivot["percent"].max()) if not pivot.empty else 100.0
    color_domain = [0.0, max(100.0, pct_max)]

    chart = (
        alt.Chart(pivot)
        .mark_rect()
        .encode(
            x=alt.X(
                "bucket_short:N",
                sort=bucket_domain,
                axis=alt.Axis(
                    title="Bucket (code)",
                    labelLimit=AXIS_TITLE_LIMIT,
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            y=alt.Y(
                "topic_graph:N",
                sort=topic_domain,
                axis=alt.Axis(
                    title="Fine topic (code)",
                    labelLimit=320,
                    labelFontSize=13,
                    labelPadding=AXIS_LABEL_PADDING,
                    labelAlign="right",
                    labelBaseline="middle",
                    labelAngle=0,
                    titleFontSize=20,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
                scale=alt.Scale(paddingInner=0.22, paddingOuter=0.06),
            ),
            color=alt.Color(
                "percent:Q",
                scale=alt.Scale(scheme="greens", domain=color_domain),
                title="% of bucket",
            ),
            tooltip=[
                alt.Tooltip("topic_bucket:N", title="Bucket"),
                alt.Tooltip("bucket_short:N", title="Code"),
                alt.Tooltip("topic_graph:N", title="Topic code"),
                alt.Tooltip("topic_names:N", title="Fine topics (full names)"),
                alt.Tooltip("percent:Q", title="% of bucket", format=".1f"),
                alt.Tooltip("count:Q", title="Sentences"),
                alt.Tooltip("bucket_total:Q", title="Bucket sentences"),
            ],
        )
        .properties(
            title="Topic mix within each bucket (% of bucket sentences)",
            width=DEFAULT_CHART_WIDTH,
            height=max(420, min(720, 120 + 48 * len(topic_domain))),
            padding={"left": 190, "right": 20, "top": 48, "bottom": 70},
        )
    )

    return chart


def compute_topic_salience(df_sent: pd.DataFrame) -> pd.DataFrame:
    """
    Per fine topic: share of its governance bucket's sentences (% of bucket).

    ``topic_bucket`` is derived from ``topic_name`` via ``TOPIC_BUCKET_MAP``, so each
    topic only appears in one bucket. "Purity" as % of topic in its top bucket is
    always 100% in that setup; salience is the comparable metric that varies and
    matches the non-zero cell per row in the topic drift heatmap.
    """
    if df_sent is None or df_sent.empty:
        return pd.DataFrame()
    if "topic_bucket" not in df_sent.columns or "topic_name" not in df_sent.columns:
        return pd.DataFrame()

    df = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df = df[df["topic_bucket"].ne("Other")]

    if df.empty:
        return pd.DataFrame()

    df["count"] = 1
    topic_counts = (
        df.groupby(["topic_name", "topic_bucket"], as_index=False)["count"]
        .sum()
        .rename(columns={"count": "sentences_in_topic"})
    )
    bucket_totals = (
        df.groupby("topic_bucket", as_index=False)["count"]
        .sum()
        .rename(columns={"count": "bucket_sentences"})
    )
    out = topic_counts.merge(bucket_totals, on="topic_bucket")
    out["pct_of_bucket"] = (out["sentences_in_topic"] / out["bucket_sentences"]) * 100.0
    return out.sort_values("pct_of_bucket", ascending=False)


def build_topic_salience_bar_chart(df_sent: pd.DataFrame, top_n: int = 8):
    salience = compute_topic_salience(df_sent)
    if salience is None or salience.empty:
        return None

    df = salience.head(top_n).copy()
    df["topic_graph"] = df["topic_name"].map(TOPIC_GRAPH_LABEL).fillna(df["topic_name"])
    known_labels = [TOPIC_GRAPH_LABEL[t] for t in TOPIC_DEFINITIONS.keys() if t in TOPIC_GRAPH_LABEL]
    present = set(df["topic_graph"].astype(str))
    topic_domain = [lab for lab in known_labels if lab in present] + sorted(
        present.difference(known_labels), key=str
    )

    bars = (
        alt.Chart(df)
        .mark_bar(color=SECONDARY_BLUE)
        .encode(
            x=alt.X(
                "pct_of_bucket:Q",
                axis=alt.Axis(
                    title="% of bucket sentences (salience)",
                    labelPadding=AXIS_LABEL_PADDING,
                    titlePadding=AXIS_TITLE_PADDING,
                    offset=8,
                    tickSize=6,
                ),
                scale=alt.Scale(domain=[0, 100]),
            ),
            y=alt.Y(
                "topic_graph:N",
                sort=topic_domain,
                axis=alt.Axis(
                    title="Fine topic (code)",
                    labelLimit=AXIS_TITLE_LIMIT,
                    labelPadding=AXIS_LABEL_PADDING,
                    labelAlign="right",
                    labelBaseline="middle",
                    labelAngle=0,
                    titlePadding=AXIS_TITLE_PADDING,
                ),
            ),
            tooltip=[
                "topic_name",
                "topic_graph",
                "topic_bucket",
                "sentences_in_topic",
                "bucket_sentences",
                "pct_of_bucket",
            ],
        )
    )

    chart = bars.properties(
        title="Fine Topic Salience",
        width=DEFAULT_CHART_WIDTH,
        height=max(DEFAULT_CHART_HEIGHT, 430),
        padding={"left": 130, "right": 20, "top": 20, "bottom": 60},
    )

    return chart


def apply_dashboard_theme_styles():
    """
    Align dashboard typography and page styling with chart design system.
    """
    st.markdown(
        f"""
        <style>
        .stApp {{
            background-color: #FFFFFF !important;
            color: {SECONDARY_BLUE};
            font-family: Garamond, "Times New Roman", serif;
        }}
        [data-testid="stSidebar"] {{
            background-color: #FFFFFF !important;
            border-right: 1px solid {GRID_COLOR};
        }}
        [data-testid="stHeader"] {{
            background: #FFFFFF;
        }}
        h1, h2, h3, .stTitle {{
            color: {PRIMARY_BLUE};
            font-family: "Helvetica Neue", Helvetica, Arial, sans-serif !important;
            letter-spacing: 0.2px;
        }}
        .stMarkdown, .stText, p, li, label, div[data-testid="stCaptionContainer"] {{
            color: {SECONDARY_BLUE};
            font-family: Garamond, "Times New Roman", serif !important;
        }}
        .stButton > button,
        .stDownloadButton > button,
        .stFormSubmitButton > button {{
            font-family: "Helvetica Neue", Helvetica, Arial, sans-serif !important;
            color: #FFFFFF !important;
            background-color: {PRIMARY_BLUE} !important;
            border: 1px solid {PRIMARY_BLUE} !important;
            border-radius: 8px !important;
            font-weight: 600 !important;
        }}
        .stButton > button *,
        .stDownloadButton > button *,
        .stFormSubmitButton > button * {{
            color: #FFFFFF !important;
        }}
        .stButton > button:hover,
        .stDownloadButton > button:hover,
        .stFormSubmitButton > button:hover {{
            background-color: {SECONDARY_BLUE} !important;
            border-color: {SECONDARY_BLUE} !important;
            color: #FFFFFF !important;
        }}
        div[data-testid="stMetricLabel"] * {{
            color: {SECONDARY_BLUE} !important;
            font-family: Garamond, "Times New Roman", serif !important;
        }}
        div[data-testid="stMetricValue"] * {{
            color: {PRIMARY_BLUE} !important;
            font-family: "Helvetica Neue", Helvetica, Arial, sans-serif !important;
            font-weight: 700 !important;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def _configured_app_password() -> str:
    """
    Read optional app password from Streamlit secrets or environment.
    If empty, password protection is disabled.
    """
    secret_value = ""
    try:
        secret_value = str(st.secrets.get("APP_PASSWORD", "")).strip()
    except Exception:
        secret_value = ""
    if secret_value:
        return secret_value
    return str(os.getenv("APP_PASSWORD", "")).strip()


def require_app_password() -> None:
    """
    Gate app access behind a single shared password when configured.
    """
    configured = _configured_app_password()
    if not configured:
        return

    if st.session_state.get("authenticated", False):
        return

    st.title("Media Intelligence Dashboard")
    st.info("This dashboard is password protected.")

    entered = st.text_input("Enter password", type="password")
    if st.button("Unlock dashboard", type="primary"):
        if hmac.compare_digest(str(entered), configured):
            st.session_state["authenticated"] = True
            st.rerun()
        else:
            st.error("Incorrect password.")
    st.stop()


def render_executive_summary_page(df_sent: pd.DataFrame):
    st.header("Executive Summary")

    df_article_sent = compute_article_sentiment(df_sent)
    st.info(
        "**Sentence-level sentiment** scores each extracted sentence with the ensemble models. "
        "**Article-level tone** aggregates those sentence scores per article. "
        "**Bucket polarity** is sentence-weighted sentiment *within* each strategic bucket (short codes CB, WO, PS, GL): "
        "for sentences assigned to a bucket, we mix their sentiment proportions using fixed numeric weights—it is not a separate polarity score for the topic label itself. "
        "**Calibrated media tone** is used only for the speedometer, with positive signals lightly upweighted to offset the cautious baseline of financial journalism."
    )

    st.markdown("### Overall Sentiment — Sentences vs Articles")
    col_a, col_b = st.columns(2)

    with col_a:
        st.markdown("Sentence Sentiment Distribution")
        chart_sent = build_sentence_distribution_chart(df_sent)
        if chart_sent is not None:
            st.altair_chart(chart_sent, use_container_width=True)
        else:
            st.write("No sentence-level sentiment data available.")

    with col_b:
        st.markdown("Article Tone Distribution")
        chart_articles = build_article_tone_chart(df_article_sent)
        if chart_articles is not None:
            st.altair_chart(chart_articles, use_container_width=True)
        else:
            st.write("No article-level sentiment data available.")

    df_polarity = compute_bucket_polarity(df_sent)
    summary_text = build_executive_summary(df_polarity)

    st.markdown("### Narrative Overview")
    st.write(summary_text)

    overall_sentence = compute_overall_score(df_sent)
    overall_article = compute_article_overall_score(df_article_sent)
    sentence_gauge_col, article_gauge_col, metric_col = st.columns([1, 1, 1])
    with sentence_gauge_col:
        st.markdown("### Calibrated Media Tone Gauge (Sentences)")
        gauge_fig = build_overall_gauge_figure(
            score=overall_sentence,
            title="Calibrated Media Tone Gauge (Sentence-level)",
            subtitle="calibrated media tone (sentences)",
        )
        st.pyplot(gauge_fig, use_container_width=False)
        plt.close(gauge_fig)
    with article_gauge_col:
        st.markdown("### Calibrated Media Tone Gauge (Articles)")
        article_gauge_fig = build_overall_gauge_figure(
            score=overall_article,
            title="Calibrated Media Tone Gauge (Article-level)",
            subtitle="calibrated media tone (articles)",
        )
        st.pyplot(article_gauge_fig, use_container_width=False)
        plt.close(article_gauge_fig)
    with metric_col:
        st.metric("Sentence Gauge Score", f"{overall_sentence:.1f} / 100")
        st.metric("Article Gauge Score", f"{overall_article:.1f} / 100")
        st.caption("Gauge calibration only: positive signals are lightly upweighted; raw charts and bucket polarity remain symmetric.")
        diag = compute_override_diagnostics(df_sent)
        st.metric("Total Overrides", f"{diag['override_count']}")
        st.caption(
            f"Sentiment overrides: {diag['sentiment_override_count']} | "
            f"Topic overrides: {diag['topic_override_count']} | "
            f"Avg override topic score: {diag['avg_override_topic_score']:.2f}"
        )

    st.markdown("### Polarity by Bucket")

    if not df_polarity.empty:
        chart_pol = build_bucket_polarity_bar_chart(df_polarity)
        st.altair_chart(chart_pol, use_container_width=True)
    else:
        st.write("No polarity data available.")

    st.markdown("### Top Sentiment Drivers (Global)")

    pos, neg = get_global_sentiment_drivers(df_sent, top_n=3)
    col1, col2 = st.columns(2)

    with col1:
        st.markdown("Top Positive Drivers")
        if not pos.empty:
            st.dataframe(
                pos[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "driver_score",
                    ]
                ]
            )
        else:
            st.write("No positive drivers available.")

    with col2:
        st.markdown("Top Negative Drivers")
        if not neg.empty:
            st.dataframe(
                neg[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "driver_score",
                    ]
                ]
            )
        else:
            st.write("No negative drivers available.")


def build_executive_summary(df_polarity: pd.DataFrame) -> str:
    if df_polarity.empty:
        return "No sentiment data available for an executive summary."

    lines = []

    for _, row in df_polarity.sort_values("polarity", ascending=False).iterrows():
        bucket = row["topic_bucket"]
        pol = row["polarity"]
        pos = row["positive_percent"]
        neg = row["negative_percent"]
        neu = row["neutral_percent"]

        if pol > 0:
            direction = "is a positive area with favourable coverage."
        elif pol < 0:
            direction = "is a negative area with critical coverage."
        else:
            direction = "shows a mixed or neutral sentiment profile."

        lines.append(
            f"{bucket} {direction} "
            f"Approx. {pos:.0f}% positive, {neg:.0f}% negative, and {neu:.0f}% neutral."
        )

    return " ".join(lines)


def render_topic_buckets_page(df_sent: pd.DataFrame, df_topics: pd.DataFrame, bucket_sizes: pd.DataFrame):
    st.header("High‑Level Topic Buckets")

    if "topic_bucket" not in df_sent.columns:
        df_sent, df_topics, bucket_sizes = apply_topic_buckets(df_sent, df_topics)

    df_sent_b = df_sent[df_sent["topic_bucket"].ne("None")].copy()
    df_sent_b = df_sent_b[df_sent_b["topic_bucket"].ne("Other")]

    st.markdown("### Override Footprint")
    diag = compute_override_diagnostics(df_sent)
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Overrides", f"{diag['override_count']}")
    c2.metric("Sentiment Overrides", f"{diag['sentiment_override_count']}")
    c3.metric("Topic Overrides", f"{diag['topic_override_count']}")
    c4.metric("Avg Topic Score (overrides)", f"{diag['avg_override_topic_score']:.2f}")

    st.markdown("### Bucket Sizes")
    if bucket_sizes is not None and not bucket_sizes.empty:
        st.dataframe(bucket_sizes[bucket_sizes["topic_bucket"].ne("None")])
    else:
        st.write("No bucket size data available.")

    df_polarity = compute_bucket_polarity(df_sent_b)

    st.markdown("### Bucket Polarity Ranking")
    if not df_polarity.empty:
        st.dataframe(df_polarity.sort_values("polarity", ascending=False))
    else:
        st.write("No polarity data available.")

    st.markdown("### Polarity by Bucket")

    if not df_polarity.empty:
        chart_pol = build_bucket_polarity_bar_chart(df_polarity)
        st.altair_chart(chart_pol, use_container_width=True)
    else:
        st.write("No polarity scores to display.")

    st.markdown("### Bucket × Sentiment View")

    view_type = st.radio(
        "Choose sentiment view",
        ["Heatmap", "Bubble chart"],
        horizontal=True,
    )

    if view_type == "Heatmap":
        heatmap = bucket_sentiment_heatmap(df_sent_b)
        if heatmap is not None:
            st.altair_chart(heatmap, use_container_width=True)
        else:
            st.write("No data available for bucket heatmap.")
    else:
        bubble = bucket_sentiment_bubble(df_sent_b)
        if bubble is not None:
            st.altair_chart(bubble, use_container_width=True)
        else:
            st.write("No data available for bucket bubble view.")

    st.markdown("### Net Balance vs Intensity (Bubble Map)")

    balance_chart = bucket_balance_bubble(df_sent_b)
    if balance_chart is not None:
        st.altair_chart(balance_chart, use_container_width=True)
    else:
        st.write("No data available for balance bubble map.")

    st.markdown("### Topic Drift (Fine Topics → Buckets)")
    st.caption(
        "Chart axes use short codes to save space. Buckets: CB = Customer & Brand Experience, "
        "GL = Governance, Leadership & Accountability, PS = Performance & Strategy, "
        "WO = Workforce, Culture & Operations. Fine topics: CR, CX, FP, LG, PO, RC, ST, WF — "
        "hover a cell for the full topic and bucket names."
    )

    drift_chart = topic_drift_heatmap(df_sent_b)
    if drift_chart is not None:
        st.altair_chart(drift_chart, use_container_width=True)
    else:
        st.write("No data available for topic drift.")

    st.markdown("### Fine-topic salience (% of bucket)")

    st.caption(
        "Each fine topic maps to one governance bucket. “Purity” as share of the topic in its "
        "strongest bucket would always read 100% here; instead we show salience — what fraction "
        "of that bucket’s sentences belong to each fine topic (same scale as the drift heatmap row)."
    )

    sal_chart = build_topic_salience_bar_chart(df_sent_b, top_n=10)
    if sal_chart is not None:
        st.altair_chart(sal_chart, use_container_width=True)

    salience_tbl = compute_topic_salience(df_sent_b)
    if not salience_tbl.empty:
        disp = salience_tbl.copy()
        disp["code"] = disp["topic_name"].map(TOPIC_GRAPH_LABEL).fillna(disp["topic_name"])
        st.dataframe(
            disp[
                [
                    "code",
                    "topic_name",
                    "topic_bucket",
                    "sentences_in_topic",
                    "bucket_sentences",
                    "pct_of_bucket",
                ]
            ].rename(
                columns={
                    "pct_of_bucket": "% of bucket",
                    "sentences_in_topic": "Sentences (topic)",
                    "bucket_sentences": "Sentences (bucket)",
                }
            )
        )
    else:
        st.write("No salience data available.")

    st.markdown("### Top Sentiment Drivers (by Bucket)")

    drivers = get_sentiment_drivers(df_sent_b)
    for bucket, d in drivers.items():
        st.markdown(f"#### {bucket}")
        col1, col2 = st.columns(2)

        with col1:
            st.markdown("Top Positive Drivers")
            if not d["positive"].empty:
                st.dataframe(d["positive"])
            else:
                st.write("No positive drivers.")

        with col2:
            st.markdown("Top Negative Drivers")
            if not d["negative"].empty:
                st.dataframe(d["negative"])
            else:
                st.write("No negative drivers.")

    st.markdown("### Narrative Summaries")

    summaries = generate_bucket_summary(df_sent_b)
    if summaries:
        for bucket, text in summaries.items():
            st.markdown(f"{bucket}")
            st.write(text)
    else:
        st.write("No summaries available.")


def sentiment_to_rgb(sentiment, topic_score):
    GREEN = (0, 153, 76)
    RED = (204, 0, 0)
    NEUTRAL = (150, 150, 150)

    try:
        val = float(topic_score)
    except Exception:
        val = None

    if val is not None:
        if val > 5:
            return GREEN
        if val < -5:
            return RED
        return NEUTRAL

    s = (sentiment or "").lower()
    if "positive" in s:
        return GREEN
    if "negative" in s:
        return RED
    return NEUTRAL


def build_sentiment_wordcloud_data(df_sent, search_terms):
    terms = [t.strip().lower() for t in search_terms.split(",") if t.strip()]
    if not terms:
        return {}, {}

    mask = df_sent["sentence"].str.lower().apply(
        lambda s: any(t in s for t in terms)
    )
    df = df_sent[mask].copy()

    if df.empty:
        return {}, {}

    freq = {t: df["sentence"].str.lower().str.count(t).sum() for t in terms}

    # Word-cloud color calibration:
    # - keep representation data-driven using mean sentiment weight per term
    # - allow a modest positive presentation shift (only after enough evidence)
    POSITIVE_SHIFT = 0.15
    MIN_MENTIONS_FOR_COLOR = 3
    POSITIVE_THRESHOLD = 0.35
    NEGATIVE_THRESHOLD = -0.35

    colors = {}
    for t in terms:
        subset = df[df["sentence"].str.lower().str.contains(t)]
        if subset.empty:
            colors[t] = (180, 180, 180)
            continue

        mention_count = int(len(subset))
        if mention_count < MIN_MENTIONS_FOR_COLOR:
            colors[t] = (150, 150, 150)
            continue

        sentiment_vals = subset["sentiment_display"].map(SENTIMENT_WEIGHTS).fillna(0.0)
        mean_sentiment = float(sentiment_vals.mean())
        calibrated = mean_sentiment + (POSITIVE_SHIFT if mean_sentiment > 0 else 0.0)

        if calibrated >= POSITIVE_THRESHOLD:
            colors[t] = (0, 153, 76)
        elif calibrated <= NEGATIVE_THRESHOLD:
            colors[t] = (204, 0, 0)
        else:
            colors[t] = (150, 150, 150)

    return freq, colors


def make_color_func(color_map):
    def color_func(word, *args, **kwargs):
        rgb = color_map.get(word, (180, 180, 180))
        return f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
    return color_func


def render_overview_page(df_sent, df_topics, df_articles):
    st.header("Overview")

    col1, col2, col3 = st.columns(3)

    with col1:
        st.metric("Total Sentences", len(df_sent))

    with col2:
        st.metric("Unique Topics", df_topics["topic_name"].nunique() if not df_topics.empty else 0)

    with col3:
        st.metric("Total Articles", len(df_articles) if df_articles is not None else 0)

    st.markdown("### Sample of Sentences")
    st.dataframe(df_sent.head(20))


def render_topic_explorer_page(df_sent):
    st.header("Topic Explorer")

    if "topic_name" not in df_sent.columns or df_sent["topic_name"].nunique() == 0:
        st.write("No topic data available.")
        return

    topic_list = sorted(df_sent["topic_name"].dropna().unique())
    topic = st.selectbox("Select a topic", topic_list)

    df_t = df_sent[df_sent["topic_name"] == topic].copy()

    st.markdown(f"### Sentences for: {topic}")
    st.dataframe(
        df_t[
            [
                "sentence",
                "topic_name",
                "topic_bucket",
                "sentiment_display",
                "topic_score",
            ]
        ]
    )


def render_entity_explorer_page(df_entities, df_sent):
    st.header("Entity Explorer")

    if df_entities.empty:
        st.write("No entity data available.")
        return

    # Fallback: check for available entity column names
    entity_col = None
    for col in ["entity_canonical", "entity", "entity_name", "text"]:
        if col in df_entities.columns:
            entity_col = col
            break
    
    if entity_col is None:
        st.write("No entity column found in data. Available columns: " + ", ".join(df_entities.columns.tolist()))
        return

    entity_list = sorted(df_entities[entity_col].dropna().unique())
    entity = st.selectbox("Select an entity", entity_list)

    df_e = df_entities[df_entities[entity_col] == entity].copy()
    st.markdown(f"### Mentions of: {entity}")
    st.dataframe(df_e)

    if "sentence_id" in df_e.columns and "global_index" in df_sent.columns:
        linked = df_sent[df_sent["global_index"].isin(df_e["sentence_id"])]
        if not linked.empty:
            st.markdown("### Sentences mentioning this entity")
            st.dataframe(
                linked[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "topic_score",
                    ]
                ]
            )
    elif "sentence_index" in df_e.columns and "global_index" in df_sent.columns:
        linked = df_sent[df_sent["global_index"].isin(df_e["sentence_index"])]
        if not linked.empty:
            st.markdown("### Sentences mentioning this entity")
            st.dataframe(
                linked[
                    [
                        "sentence",
                        "topic_name",
                        "topic_bucket",
                        "sentiment_display",
                        "topic_score",
                    ]
                ]
            )


def _ppt_image_pixel_size(path: Path) -> tuple:
    if Image is None:
        return (1200, 700)
    try:
        with Image.open(path) as im:
            return int(im.size[0]), int(im.size[1])
    except Exception:
        return (1200, 700)


def _ppt_add_picture_fit(slide, image_path: str, left_in: float, top_in: float, max_w_in: float, max_h_in: float):
    """
    Insert a picture scaled to fit max_w_in x max_h_in while preserving aspect ratio.
    Sets lock aspect ratio when supported so manual resize in PowerPoint stays proportional.
    """
    p = Path(image_path)
    if not p.exists():
        return None

    px_w, px_h = _ppt_image_pixel_size(p)
    if px_w <= 0 or px_h <= 0:
        px_w, px_h = 1200, 700

    scale = min(max_w_in / px_w, max_h_in / px_h)
    w_in = max(0.01, px_w * scale)
    h_in = max(0.01, px_h * scale)

    pic = slide.shapes.add_picture(str(p), Inches(left_in), Inches(top_in), width=Inches(w_in), height=Inches(h_in))
    _ppt_lock_picture_aspect_ratio(pic)
    return pic


def _ppt_lock_picture_aspect_ratio(pic):
    """Lock inserted PNGs so PowerPoint preserves proportions during resize/move."""
    try:
        pic.lock_aspect_ratio = True
    except Exception:
        pass

    try:
        cnv_pic_pr = pic._element.xpath(".//p:cNvPicPr")[0]
        pic_locks = cnv_pic_pr.find(qn("a:picLocks"))
        if pic_locks is None:
            pic_locks = OxmlElement("a:picLocks")
            cnv_pic_pr.append(pic_locks)
        pic_locks.set("noChangeAspect", "1")
    except Exception:
        pass


def _ppt_add_title(slide, title_text: str):
    title_shape = slide.shapes.title
    title_shape.text = title_text
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = "Helvetica Neue"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PPT_PRIMARY_RGB
    p.alignment = PP_ALIGN.LEFT


def _ppt_add_subtitle(slide, subtitle_text: str):
    tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.2), Inches(0.45))
    tf = tx_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.name = "Garamond"
    p.font.size = Pt(16)
    p.font.color.rgb = PPT_SECONDARY_RGB


def _ppt_add_footer_key(slide):
    footer_box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.45), Inches(0.28))
    tf = footer_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = PPT_FOOTER_KEY_TEXT
    p.font.name = "Garamond"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(95, 95, 95)
    p.alignment = PP_ALIGN.LEFT


def _ppt_add_bullets(
    slide,
    bullets: list,
    left: float = 0.8,
    top: float = 1.6,
    width: float = 11.6,
    height: float = 4.8,
    font_pt: int = 18,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Garamond"
        p.font.size = Pt(int(font_pt))
        p.font.color.rgb = PPT_SECONDARY_RGB


def _ppt_place_images(slide, image_paths: list):
    """Place one or more chart PNGs on the right side of a slide with consistent spacing."""
    images = [p for p in (image_paths or []) if p and Path(str(p)).exists()]
    if not images:
        return

    chart_left = 6.35
    chart_top = 1.25
    max_w = 6.75
    panel_h = 5.85
    gap = 0.16

    n = len(images)
    if n == 1:
        _ppt_add_picture_fit(slide, images[0], chart_left, chart_top, max_w, panel_h)
        return

    if n == 2:
        slot_h = (panel_h - gap) / 2.0
        _ppt_add_picture_fit(slide, images[0], chart_left, chart_top, max_w, slot_h)
        _ppt_add_picture_fit(slide, images[1], chart_left, chart_top + slot_h + gap, max_w, slot_h)
        return

    if n == 3:
        top_h = (panel_h - gap) * 0.58
        bottom_h = panel_h - top_h - gap
        col_w = (max_w - gap) / 2.0
        _ppt_add_picture_fit(slide, images[0], chart_left, chart_top, col_w, top_h)
        _ppt_add_picture_fit(slide, images[1], chart_left + col_w + gap, chart_top, col_w, top_h)
        _ppt_add_picture_fit(slide, images[2], chart_left, chart_top + top_h + gap, max_w, bottom_h)
        return

    # 4+ images: 2x2 grid (use first four)
    row_h = (panel_h - gap) / 2.0
    col_w = (max_w - gap) / 2.0
    positions = [
        (chart_left, chart_top, col_w, row_h),
        (chart_left + col_w + gap, chart_top, col_w, row_h),
        (chart_left, chart_top + row_h + gap, col_w, row_h),
        (chart_left + col_w + gap, chart_top + row_h + gap, col_w, row_h),
    ]
    for i in range(min(4, len(images))):
        left, top, w, h = positions[i]
        _ppt_add_picture_fit(slide, images[i], left, top, w, h)


def _safe_str(value, default: str = "n/a") -> str:
    if value is None:
        return default
    txt = str(value).strip()
    return txt if txt else default


def render_short_code_key():
    st.markdown("---")
    st.caption("**Short-code key (used across charts)**")
    st.caption(BUCKET_KEY_TEXT)
    st.caption(FINE_TOPIC_KEY_TEXT)


def compute_overall_score(df_sent: pd.DataFrame) -> float:
    if df_sent is None or df_sent.empty:
        return 50.0
    weights = df_sent["sentiment_display"].map(GAUGE_SENTIMENT_WEIGHTS).fillna(0.0)
    weighted_mean = float(weights.mean())
    score_0_100 = max(0.0, min(100.0, ((weighted_mean + 2.0) / 4.0) * 100.0))
    return score_0_100


def compute_article_overall_score(df_article_sent: pd.DataFrame) -> float:
    if df_article_sent is None or df_article_sent.empty or "avg_weight" not in df_article_sent.columns:
        return 50.0
    avg_weights = pd.to_numeric(df_article_sent["avg_weight"], errors="coerce").dropna()
    if avg_weights.empty:
        return 50.0
    calibrated = avg_weights.apply(lambda x: x * 1.25 if x > 0 else x).clip(lower=-2.0, upper=2.25)
    weighted_mean = float(calibrated.mean())
    score_0_100 = max(0.0, min(100.0, ((weighted_mean + 2.0) / 4.0) * 100.0))
    return score_0_100


def build_overall_gauge_figure(score: float, title: str = "Calibrated Media Tone Gauge", subtitle: str = "calibrated media tone"):
    fig, ax = plt.subplots(figsize=(5.2, 2.8))
    fig.patch.set_facecolor("none")
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    ax.patch.set_alpha(0)
    ax.set_title(title, color=PRIMARY_BLUE, fontsize=14, fontweight="bold", pad=12)

    bands = [(0, 40, "#d73027"), (40, 60, "#fdae61"), (60, 100, "#1a9850")]
    for s, e, color in bands:
        a1 = 180 - (s / 100) * 180
        a2 = 180 - (e / 100) * 180
        wedge = plt.matplotlib.patches.Wedge((0, 0), 1.0, a2, a1, width=0.28, facecolor=color, edgecolor="white")
        ax.add_patch(wedge)

    angle = np.pi * (1 - score / 100.0)
    x, y = 0.78 * np.cos(angle), 0.78 * np.sin(angle)
    ax.plot([0, x], [0, y], color=PRIMARY_BLUE, linewidth=3)
    ax.scatter([0], [0], color=PRIMARY_BLUE, s=35, zorder=5)
    ax.text(0, -0.12, f"{score:.1f} / 100", ha="center", va="center", color=PRIMARY_BLUE, fontsize=14)
    ax.text(0, -0.26, subtitle, ha="center", va="center", color=SECONDARY_BLUE, fontsize=10)

    ax.set_xlim(-1.1, 1.1)
    ax.set_ylim(-0.32, 1.1)
    ax.axis("off")
    return fig


def compute_override_diagnostics(df_sent: pd.DataFrame) -> dict:
    if df_sent is None or df_sent.empty:
        return {
            "override_count": 0,
            "sentiment_override_count": 0,
            "topic_override_count": 0,
            "avg_override_topic_score": 0.0,
        }

    sent_override_col = "manual_sentiment_override_applied"
    topic_override_col = "manual_topic_override_applied"
    sent_overrides = int(df_sent[sent_override_col].fillna(False).sum()) if sent_override_col in df_sent.columns else 0
    topic_overrides = int(df_sent[topic_override_col].fillna(False).sum()) if topic_override_col in df_sent.columns else 0
    override_count = int(
        (df_sent.get(sent_override_col, pd.Series(False, index=df_sent.index)).fillna(False) |
         df_sent.get(topic_override_col, pd.Series(False, index=df_sent.index)).fillna(False)).sum()
    )

    avg_score = 0.0
    if topic_overrides > 0 and "topic_score" in df_sent.columns and topic_override_col in df_sent.columns:
        avg_score = float(df_sent[df_sent[topic_override_col].fillna(False)]["topic_score"].fillna(0.0).mean())

    return {
        "override_count": override_count,
        "sentiment_override_count": sent_overrides,
        "topic_override_count": topic_overrides,
        "avg_override_topic_score": avg_score,
    }


def _polarity_ranking_lines(df_polarity: pd.DataFrame, limit: int = 4) -> list:
    if df_polarity is None or df_polarity.empty:
        return ["No polarity ranking available."]
    lines = []
    for _, row in df_polarity.sort_values("polarity", ascending=False).head(limit).iterrows():
        lines.append(
            f"{_safe_str(row['topic_bucket'])}: polarity {float(row['polarity']):.1f} | "
            f"+{float(row['positive_percent']):.1f}% / -{float(row['negative_percent']):.1f}%"
        )
    return lines


def _storyboard_taxonomy_slide_payload() -> dict:
    """First PPTX slide: full strategic bucket names and fine topic names with definitional anchors."""
    from collections import defaultdict

    bucket_to_topics: dict[str, list] = defaultdict(list)
    fine_topics = list(TOPIC_DEFINITIONS.keys())
    for t in fine_topics:
        b = TOPIC_BUCKET_MAP.get(t, "None")
        if b not in ("None", "Other"):
            bucket_to_topics[b].append(t)

    left_lines = [
        "Strategic governance buckets (short codes used on charts):",
    ]
    for b in BUCKET_ORDER:
        code = BUCKET_SHORT.get(b, "")
        subs = "; ".join(sorted(bucket_to_topics.get(b, []), key=str.lower))
        left_lines.append(f"{code} — {b}")
        left_lines.append(f"Fine topics: {subs}")

    right_lines = [
        "Fine-grained classifier topics (full names; axis codes in parentheses):",
    ]
    for t in sorted(fine_topics, key=str.lower):
        code = TOPIC_GRAPH_LABEL.get(t, "—")
        b = TOPIC_BUCKET_MAP.get(t, "None")
        bcode = BUCKET_SHORT.get(b, b)
        anchors = TOPIC_DEFINITIONS.get(t) or []
        first = str(anchors[0]) if anchors else ""
        if len(first) > 118:
            first = first[:115] + "…"
        right_lines.append(f"{t} ({code} → {bcode})")
        right_lines.append(f"Anchor: {first}")

    return {
        "title": "Topic & bucket definitions",
        "subtitle": "Full names for governance buckets and classifier topics (first definitional anchor each)",
        "layout": "two_column_text",
        "left_bullets": left_lines,
        "right_bullets": right_lines,
        "images": [],
        "bullet_font_pt": 11,
    }


def build_storyboard_slides(df_sent: pd.DataFrame, df_article_sent: pd.DataFrame, df_polarity: pd.DataFrame, bucket_sizes: pd.DataFrame):
    slides = []
    slides.append(_storyboard_taxonomy_slide_payload())
    overall_score = compute_overall_score(df_sent)
    overall_article_score = compute_article_overall_score(df_article_sent)
    pos, neg = get_global_sentiment_drivers(df_sent, top_n=3)
    diag = compute_override_diagnostics(df_sent)
    bucket_summaries = generate_bucket_summary(df_sent)

    df_plot = df_sent[df_sent.get("topic_bucket", "None").ne("None")].copy() if "topic_bucket" in df_sent.columns else df_sent.copy()
    if "topic_bucket" in df_plot.columns:
        df_plot = df_plot[df_plot["topic_bucket"].ne("Other")]

    chart_sentence_path = export_sentence_distribution_png(df_sent, "slide_sentence_distribution.png")
    chart_article_path = export_article_tone_png(df_article_sent, "slide_article_distribution.png")
    chart_balance_path = export_bucket_balance_png(df_sent, "slide_balance_map.png")
    chart_polarity_bar_path = export_bucket_polarity_png(df_polarity, "slide_bucket_polarity_bar.png")
    chart_bucket_sizes_path = export_bucket_sizes_png(bucket_sizes, "slide_bucket_sizes.png")
    chart_bucket_heat_path = export_bucket_sentiment_heatmap_png(df_plot, "slide_bucket_sentiment_heatmap.png")
    chart_bucket_bubble_path = export_bucket_sentiment_bubble_png(df_plot, "slide_bucket_sentiment_bubble.png")
    chart_drift_path = export_topic_drift_heatmap_png(df_plot, "slide_topic_drift.png")
    chart_salience_path = export_topic_salience_png(df_plot, "slide_topic_salience.png", top_n=8)

    gauge_fig = build_overall_gauge_figure(
        score=overall_score,
        title="Calibrated Media Tone Gauge (Sentence-level)",
        subtitle="calibrated media tone (sentences)",
    )
    gauge_path = CHART_EXPORT_DIR / "slide_gauge_sentence.png"
    gauge_fig.savefig(gauge_path, dpi=170, bbox_inches="tight", facecolor="none", edgecolor="none", transparent=True)
    plt.close(gauge_fig)

    gauge_article_fig = build_overall_gauge_figure(
        score=overall_article_score,
        title="Calibrated Media Tone Gauge (Article-level)",
        subtitle="calibrated media tone (articles)",
    )
    gauge_article_path = CHART_EXPORT_DIR / "slide_gauge_article.png"
    gauge_article_fig.savefig(gauge_article_path, dpi=170, bbox_inches="tight", facecolor="none", edgecolor="none", transparent=True)
    plt.close(gauge_article_fig)

    slides.append({
        "title": "Executive Overview",
        "subtitle": "Sentiment and topic governance snapshot",
        "bullets": [
            f"Calibrated media tone (sentences): {overall_score:.1f} / 100",
            f"Calibrated media tone (articles): {overall_article_score:.1f} / 100",
            f"Analysed sentences: {len(df_sent):,}",
            f"Analysed articles: {df_article_sent['article_id'].nunique() if not df_article_sent.empty else 0:,}",
            "Coverage is classified using 8 fine topics mapped to 4 governance buckets.",
            "Speedometer applies a modest positive calibration for cautious financial-news baselines; raw distributions and bucket polarity remain symmetric.",
            "Sentence-level sentiment drives polarity and intensity. Article-level tone aggregates sentence signals.",
        ],
        "images": [
            str(p)
            for p in [gauge_path, gauge_article_path]
            if p.exists()
        ],
    })

    slides.append({
        "title": "Methodology (Readability)",
        "subtitle": "How to read the charts in this pack",
        "bullets": [
            "Slide 1 lists every strategic bucket and fine topic with full names and definitional anchors.",
            "Bucket codes CB, GL, PS, WO match slide 1; fine-topic axis codes (CR, CX, FP, LG, PO, RC, ST, WF) are shorthand for those full labels.",
            "Topic drift heatmap shows each fine topic’s share of sentences within a bucket (% of bucket). "
            "The companion bar chart is bucket salience (not purity): fine topics map 1:1 to buckets, so classic purity would always read 100%.",
            "Polarity scores are computed at sentence level within each bucket, then summarised for reporting.",
            "The speedometer is calibrated only at presentation level: Positive = +1.25, Very Positive = +2.25; negative and neutral weights stay unchanged.",
        ],
        "images": [],
    })

    slides.append({
        "title": "Sentiment Distribution",
        "subtitle": "Sentence-level and article-level sentiment composition",
        "bullets": [
            "Sentence-level bars show sentiment intensity across all extracted claims.",
            "Article tone captures aggregate media framing per article.",
        ],
        "images": [p for p in [chart_sentence_path, chart_article_path] if p and Path(p).exists()],
    })

    slides.append({
        "title": "Bucket Coverage & Polarity",
        "subtitle": "Volume and directional pressure by bucket",
        "bullets": [
            "Bucket sizes show where coverage concentrates.",
            "Polarity bars summarise net positive vs negative sentence mix by bucket.",
        ],
        "images": [p for p in [chart_bucket_sizes_path, chart_polarity_bar_path] if p and Path(p).exists()],
    })

    slides.append({
        "title": "Bucket Polarity — Narrative Ranking",
        "subtitle": "High-level bucket ranking by net polarity",
        "bullets": _polarity_ranking_lines(df_polarity, limit=8),
        "images": [chart_balance_path] if chart_balance_path and Path(chart_balance_path).exists() else [],
    })

    slides.append({
        "title": "Bucket × Sentiment Composition",
        "subtitle": "Where sentiment concentrates inside each bucket",
        "bullets": [
            "Heatmap emphasises dominant sentiment cells within each bucket.",
            "Bubble view highlights relative intensity of sentiment states.",
        ],
        "images": [p for p in [chart_bucket_heat_path, chart_bucket_bubble_path] if p and Path(p).exists()],
    })

    slides.append({
        "title": "Topic drift & bucket salience",
        "subtitle": "Fine-topic concentration inside each governance bucket",
        "bullets": [
            "Heatmap: each cell is that fine topic’s share of all sentences in the column bucket (% of bucket).",
            "Bar chart: same metric for the top fine topics — how much of each bucket they absorb, not “purity”.",
        ],
        "images": [p for p in [chart_drift_path, chart_salience_path] if p and Path(p).exists()],
    })

    bucket_lines = []
    if bucket_summaries:
        for bucket, text in list(bucket_summaries.items())[:4]:
            bucket_lines.append(f"{bucket}: {text}")
    else:
        bucket_lines.append("No bucket narrative summaries available for this export.")

    slides.append({
        "title": "Bucket Narrative Summaries",
        "subtitle": "Plain-language interpretation by bucket",
        "bullets": bucket_lines[:6],
        "images": [],
    })

    driver_bullets = []
    if not neg.empty:
        driver_bullets.append("Top negative drivers:")
        for _, row in neg.head(3).iterrows():
            driver_bullets.append(f"- [{_safe_str(row['topic_bucket'])}] {_safe_str(row['sentence'])[:140]}")
    if not pos.empty:
        driver_bullets.append("Top positive drivers:")
        for _, row in pos.head(3).iterrows():
            driver_bullets.append(f"- [{_safe_str(row['topic_bucket'])}] {_safe_str(row['sentence'])[:140]}")
    if not driver_bullets:
        driver_bullets = ["No standout global drivers identified in this export."]

    slides.append({
        "title": "Global Sentiment Drivers",
        "subtitle": "Representative sentences shaping the overall story",
        "bullets": driver_bullets[:10],
        "images": [],
    })

    action_lines = []
    if not neg.empty:
        for _, row in neg.head(3).iterrows():
            action_lines.append(f"Mitigate: {_safe_str(row['topic_bucket'])} - {_safe_str(row['sentence'])[:110]}")
    if not pos.empty:
        for _, row in pos.head(2).iterrows():
            action_lines.append(f"Amplify: {_safe_str(row['topic_bucket'])} - {_safe_str(row['sentence'])[:110]}")
    if not action_lines:
        action_lines = ["No clear action drivers were identified in the current dataset."]

    slides.append({
        "title": "Action Plan Narrative",
        "subtitle": "Priority interventions and messaging actions",
        "bullets": action_lines[:8],
        "images": [str(WORDCLOUD_EXPORT_PATH)] if WORDCLOUD_EXPORT_PATH.exists() else [],
    })

    slides.append({
        "title": "Governance & Override Review",
        "subtitle": "Classification quality controls and review workload",
        "bullets": [
            f"Total overrides applied: {diag['override_count']}",
            f"Sentiment overrides: {diag['sentiment_override_count']}",
            f"Topic overrides: {diag['topic_override_count']}",
            f"Average topic score on overrides: {diag['avg_override_topic_score']:.2f}",
            "Outlier queue prioritises low-confidence and drift-risk sentences for human review.",
            "Use the Outlier Review Queue to apply corrections; the dashboard reloads from `master.json`. Rerun the pipeline for new articles or full model recomputation.",
        ],
        "images": [],
    })

    return slides


def export_storyboard_to_pptx(slides: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[5]

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        # Clean white background panel for consistent corporate look.
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(7.5),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        _ppt_add_title(slide, s.get("title", f"Slide {idx + 1}"))
        _ppt_add_subtitle(slide, s.get("subtitle", ""))

        images = [p for p in s.get("images", []) if p and Path(str(p)).exists()]
        bullets = s.get("bullets", []) or []
        font_pt = int(s.get("bullet_font_pt", 18))

        if s.get("layout") == "two_column_text":
            _ppt_add_bullets(
                slide,
                s.get("left_bullets") or [],
                left=0.48,
                top=1.42,
                width=6.15,
                height=5.78,
                font_pt=font_pt,
            )
            _ppt_add_bullets(
                slide,
                s.get("right_bullets") or [],
                left=6.78,
                top=1.42,
                width=6.05,
                height=5.78,
                font_pt=font_pt,
            )
        elif images:
            _ppt_add_bullets(slide, bullets, left=0.65, top=1.45, width=5.9, height=5.45, font_pt=font_pt)
            _ppt_place_images(slide, images)
        else:
            _ppt_add_bullets(slide, bullets, left=0.75, top=1.45, width=11.8, height=5.45, font_pt=font_pt)
        _ppt_add_footer_key(slide)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def build_narrative(df_sent: pd.DataFrame) -> str:
    df_article_sent = compute_article_sentiment(df_sent)
    df_polarity = compute_bucket_polarity(df_sent)
    bucket_summaries = generate_bucket_summary(df_sent)
    pos, neg = get_global_sentiment_drivers(df_sent, top_n=5)
    score = compute_overall_score(df_sent)
    diag = compute_override_diagnostics(df_sent)

    lines = []
    lines.append("MEDIA INTELLIGENCE NARRATIVE")
    lines.append("=" * 30)
    lines.append("")
    lines.append(f"Overall Sentiment Score: {score:.1f}/100")
    lines.append(f"Sentence Volume: {len(df_sent):,}")
    lines.append(f"Article Volume: {df_article_sent['article_id'].nunique() if not df_article_sent.empty else 0:,}")
    lines.append(
        f"Override Footprint: {diag['override_count']} total "
        f"(sentiment {diag['sentiment_override_count']}, topic {diag['topic_override_count']}, "
        f"avg topic score {diag['avg_override_topic_score']:.2f})"
    )
    lines.append("")
    lines.append("Polarity Rankings")
    lines.append("-" * 30)
    for line in _polarity_ranking_lines(df_polarity, limit=10):
        lines.append(f"- {line}")

    lines.append("")
    lines.append("Bucket Summaries")
    lines.append("-" * 30)
    if bucket_summaries:
        for bucket, summary in bucket_summaries.items():
            lines.append(f"- {bucket}: {summary}")
    else:
        lines.append("- No bucket summary available.")

    lines.append("")
    lines.append("Primary Drivers")
    lines.append("-" * 30)
    if not neg.empty:
        lines.append("Risk Drivers:")
        for _, row in neg.iterrows():
            lines.append(f"- [{_safe_str(row['topic_bucket'])}] {_safe_str(row['sentence'])}")
    if not pos.empty:
        lines.append("Positive Drivers:")
        for _, row in pos.iterrows():
            lines.append(f"- [{_safe_str(row['topic_bucket'])}] {_safe_str(row['sentence'])}")
    if neg.empty and pos.empty:
        lines.append("- No sentiment drivers identified.")

    lines.append("")
    lines.append("Action Recommendations")
    lines.append("-" * 30)
    lines.append("- Prioritise rapid-response messaging in negatively ranked buckets.")
    lines.append("- Institutionalise sentence-level override reviews for flagged low-confidence items.")
    lines.append("- Track topic drift and bucket salience weekly to detect taxonomy slippage.")
    lines.append("- Amplify strongest positive drivers in executive and external communications.")

    return "\n".join(lines)


def main():
    st.set_page_config(page_title="Media Intelligence Dashboard", layout="wide")
    require_app_password()
    apply_dashboard_theme_styles()
    st.title("Media Intelligence Dashboard")
    st.caption("Professional sentiment intelligence, topic governance, and storyboard export")

    if not MASTER_JSON.exists():
        st.error("`master.json` not found. Please run the pipeline first.")
        return

    try:
        _, df_sent, df_topics, df_entities, _, _, _ = load_master()
    except Exception as exc:
        st.error(f"Failed to load source data: {exc}")
        return

    df_sent, df_topics, bucket_sizes = apply_topic_buckets(df_sent, df_topics)
    df_articles = compute_article_sentiment(df_sent)

    st.sidebar.header("Navigation")
    page = st.sidebar.radio(
        "Choose page",
        [
            "Overview",
            "Executive Summary",
            "Topic Buckets",
            "Outlier Review Queue",
            "Topic Explorer",
            "Entity Explorer",
            "Sentence Inspector",
            "PowerPoint Storyboard",
            "Narrative Export",
        ],
    )

    if page == "Overview":
        render_overview_page(df_sent, df_topics, df_articles)
    elif page == "Executive Summary":
        render_executive_summary_page(df_sent)
    elif page == "Topic Buckets":
        render_topic_buckets_page(df_sent, df_topics, bucket_sizes)
    elif page == "Outlier Review Queue":
        render_outlier_review_page(df_sent)
    elif page == "Topic Explorer":
        render_topic_explorer_page(df_sent)
    elif page == "Entity Explorer":
        render_entity_explorer_page(df_entities, df_sent)
    elif page == "Sentence Inspector":
        render_sentence_inspector_page(df_sent)
    elif page == "PowerPoint Storyboard":
        render_powerpoint_storyboard(df_sent, df_topics, bucket_sizes)
    elif page == "Narrative Export":
        render_narrative_export(df_sent)

    render_short_code_key()


if __name__ == "__main__":
    main()
